"""
XAI-ITD-DLP Framework - Flask Backend
Module 1: User Authentication + Context Monitoring (Complete)
Tasks:
  1. Location-Based Login and Role Authentication
  2. Real-Time Employee Activity Monitoring
  3. Security Policy Enforcement and Event Logging
"""

import random
import subprocess
import smtplib
import secrets
import requests
import threading
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime, timedelta
from functools import wraps

from flask import Flask, render_template, request, jsonify, session, redirect, url_for
from flask_socketio import SocketIO, emit
import logging
logging.getLogger('werkzeug').setLevel(logging.ERROR)
import sys, os
sys.path.insert(0, os.path.dirname(__file__))

from config import (SMTP_SERVER, SMTP_PORT, SMTP_EMAIL, SMTP_PASSWORD,
                    SECRET_KEY, OTP_EXPIRY_SECONDS)
from models.files import (
    save_file_record, get_files_for_employee, get_all_files,
    get_file_by_id, get_file_by_id_unrestricted, delete_file_record,
    create_approval_request, get_pending_approvals, get_all_approvals,
    get_employee_requests, resolve_approval, get_approval_by_id, UPLOAD_FOLDER,
    save_file_to_gridfs, get_file_from_gridfs, delete_file_from_gridfs
)
from werkzeug.utils import secure_filename
import uuid

ALLOWED_EXTENSIONS = {
    'pdf','docx','xlsx','xls','csv','txt','pptx','png','jpg','jpeg','gif','mp4','zip','py','json'
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

from models.user import (
    get_user_by_email, verify_password, save_otp, get_otp, delete_otp,
    log_activity, log_security_event, get_all_logs, get_all_security_events,
    get_user_logs, seed_users, users_col, update_user_login, normalize_existing_emails,
    is_location_allowed, set_allowed_locations, get_login_history,
    get_all_employees, create_user, get_all_users_by_role,
    get_all_logs_for_user, get_all_security_events_for_user,
    get_all_logs_unrestricted, get_all_security_events_unrestricted,
    deactivate_user, reactivate_user, get_system_stats, seed_admin,
    check_mongo_connection,
    # Device fingerprint
    save_device_profile, get_device_profile, check_device_mismatch,
    # Impossible travel
    update_login_coords, check_impossible_travel,
    # Travel mode
    request_travel_mode, approve_travel_mode, reject_travel_mode,
    get_travel_mode, get_travel_history, is_travel_mode_active, get_all_travel_requests,
)

app = Flask(__name__)
app.secret_key = SECRET_KEY
socketio = SocketIO(app, cors_allowed_origins="*", async_mode="threading", manage_session=False, logger=False, engineio_logger=False)

# ── Active session tokens — MongoDB-backed so they survive Render restarts ──
from models.user import users_col as _users_col_ref
_tokens_col = _users_col_ref.database["active_tokens"]
try:
    _tokens_col.create_index("token", unique=True)
except Exception:
    pass
try:
    _tokens_col.create_index("expires_at", expireAfterSeconds=0)
except Exception:
    pass


class _ActiveTokens:
    """dict-like {token: email} backed by MongoDB. Survives server restarts."""

    def __init__(self):
        self._cache = {}
        for doc in _tokens_col.find({}, {"_id": 0, "token": 1, "email": 1}):
            self._cache[doc["token"]] = doc["email"]

    def __contains__(self, token):
        if token in self._cache:
            return True
        doc = _tokens_col.find_one({"token": token})
        if doc:
            self._cache[token] = doc["email"]
            return True
        return False

    def __getitem__(self, token):
        if token not in self._cache:
            doc = _tokens_col.find_one({"token": token})
            if doc:
                self._cache[token] = doc["email"]
            else:
                raise KeyError(token)
        return self._cache[token]

    def __setitem__(self, token, email):
        self._cache[token] = email
        try:
            _tokens_col.update_one(
                {"token": token},
                {"$set": {"token": token, "email": email,
                           "expires_at": datetime.utcnow() + timedelta(hours=24)}},
                upsert=True
            )
        except Exception:
            pass

    def __delitem__(self, token):
        self._cache.pop(token, None)
        try:
            _tokens_col.delete_one({"token": token})
        except Exception:
            pass

    def get(self, token, default=None):
        try:
            return self[token]
        except KeyError:
            return default

    def __len__(self):
        return len(self._cache)

    def values(self):
        return self._cache.values()

    def keys(self):
        return self._cache.keys()

    def items(self):
        return self._cache.items()


active_tokens = _ActiveTokens()

# ── Meeting & Client blueprints ───────────────────────────────
from routes.meetings           import meetings_bp
from routes.meeting_attendance import attendance_bp
from routes.meeting_notes      import notes_bp, polls_bp, actions_bp, feedback_bp
from routes.client_meetings    import client_bp
from sockets_meeting           import register_meeting_sockets

app.register_blueprint(meetings_bp)
app.register_blueprint(attendance_bp)
app.register_blueprint(notes_bp)
app.register_blueprint(polls_bp)
app.register_blueprint(actions_bp)
app.register_blueprint(feedback_bp)
app.register_blueprint(client_bp)

from routes.threat_api import threat_bp
app.register_blueprint(threat_bp)
from routes.xai_api import xai_bp, init_xai_bp
app.register_blueprint(xai_bp)

# Real-time AI scorer
try:
    from ml.realtime_scorer import score_now as ai_score_now
except Exception as _e:
    ai_score_now = None
    print("[AI] realtime_scorer not available: " + str(_e))
init_xai_bp()
from models.user import db as _meeting_db
register_meeting_sockets(socketio, _meeting_db)



# --- Helpers ------------------------------------------------------------------

def _get_real_public_ip():
    """When Flask sees 127.0.0.1, fetch the machine's actual public IP."""
    for url in ["https://api.ipify.org", "https://checkip.amazonaws.com"]:
        try:
            r = requests.get(url, timeout=3)
            return r.text.strip()
        except Exception:
            continue
    return None


def _lookup_ip_location(ip):
    """Try multiple geo APIs in order - use first that returns a real city."""
    # 1. ip-api.com - best coverage for Indian ISPs, free, no key needed
    try:
        r = requests.get(
            "http://ip-api.com/json/" + ip + "?fields=status,city,regionName,country,lat,lon,org",
            timeout=4
        )
        d = r.json()
        if d.get("status") == "success" and d.get("city") and d.get("city") != "Unknown":
            return {
                "city":    d.get("city", "Unknown"),
                "region":  d.get("regionName", ""),
                "country": d.get("country", "Unknown"),
                "ip": ip, "lat": d.get("lat"), "lon": d.get("lon"), "org": d.get("org", "")
            }
    except Exception:
        pass

    # 2. ipwho.is - backup
    try:
        r = requests.get("https://ipwho.is/" + ip, timeout=4)
        d = r.json()
        if d.get("success") and d.get("city") and d.get("city") != "Unknown":
            return {
                "city":    d.get("city", "Unknown"),
                "region":  d.get("region", ""),
                "country": d.get("country", "Unknown"),
                "ip": ip, "lat": d.get("latitude"), "lon": d.get("longitude"),
                "org": d.get("connection", {}).get("org", "")
            }
    except Exception:
        pass

    # 3. ipapi.co - fallback
    try:
        r = requests.get("https://ipapi.co/" + ip + "/json/", timeout=4)
        d = r.json()
        if d.get("city") and d.get("city") != "Unknown":
            return {
                "city":    d.get("city", "Unknown"),
                "region":  d.get("region", ""),
                "country": d.get("country_name", "Unknown"),
                "ip": ip, "lat": d.get("latitude"), "lon": d.get("longitude"), "org": d.get("org", "")
            }
    except Exception:
        pass

    return {"city": "Unknown", "region": "", "country": "Unknown", "ip": ip, "lat": None, "lon": None, "org": ""}


def get_location_from_ip(ip):
    # When running locally Flask sees 127.0.0.1 - resolve the real public IP first
    if ip in ("127.0.0.1", "::1", "localhost"):
        real_ip = _get_real_public_ip()
        if real_ip:
            print("[LOCATION] Localhost -> real IP: " + real_ip)
            ip = real_ip
        else:
            print("[LOCATION] Could not resolve real public IP")
            return {"city": "Unknown", "region": "", "country": "Unknown", "ip": "127.0.0.1", "lat": None, "lon": None, "org": ""}
    return _lookup_ip_location(ip)


def send_otp_email(to_email, otp, name):
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = " XAI-ITD-DLP Login OTP"
        msg["From"] = SMTP_EMAIL
        msg["To"] = to_email
        html = """
        <div style="font-family:monospace;background:#0a0a0f;color:#00ff88;padding:30px;border-radius:10px;max-width:500px">
          <h2 style="color:#00ff88;letter-spacing:3px;">XAI-ITD-DLP SYSTEM</h2>
          <p>Hello <strong>{name}</strong>,</p>
          <p>Your one-time login code:</p>
          <div style="font-size:40px;font-weight:bold;letter-spacing:10px;color:#fff;
                      background:#111;padding:20px;border-radius:8px;text-align:center;
                      border:2px solid #00ff88;">{otp}</div>
          <p style="color:#888;margin-top:20px;"> Expires in 2 minutes. Do not share this code.</p>
          <p style="color:#ff4444;">If you did not request this, contact your administrator immediately.</p>
        </div>
        """.format(name=name, otp=otp)
        msg.attach(MIMEText(html, "html"))
        with smtplib.SMTP(SMTP_SERVER, SMTP_PORT) as server:
            server.starttls()
            server.login(SMTP_EMAIL, SMTP_PASSWORD)
            server.sendmail(SMTP_EMAIL, to_email, msg.as_string())
        print("[EMAIL] OTP sent successfully to " + to_email)
        return True
    except Exception as e:
        print("[EMAIL ERROR] " + str(e))
        return False


# ─────────────────────────────────────────────────────────────────────────────
# MEETING INVITE EMAIL
# ─────────────────────────────────────────────────────────────────────────────
# ↓↓↓ PASTE YOUR CLIENT EMAIL ADDRESS HERE ↓↓↓
CLIENT_MEETING_EMAIL = "robinescobar220@gmail.com"
# ↑↑↑ Replace the above with the actual client email address ↑↑↑


def send_meeting_invite_email(to_email, client_name, employee_name,
                              join_url, password, scheduled_time, agenda):
    """
    Send meeting join link + password to recipient email.
    Uses same SMTP config as OTP emails.
    Supports port 587 (STARTTLS) and port 465 (SSL).
    """
    print("[EMAIL] Sending meeting invite to: " + str(to_email))
    print("[EMAIL] SMTP: " + str(SMTP_SERVER) + ":" + str(SMTP_PORT) + " from " + str(SMTP_EMAIL))
    try:
        msg            = MIMEMultipart("alternative")
        msg["Subject"] = "Your Meeting Invitation"
        msg["From"]    = SMTP_EMAIL
        msg["To"]      = to_email

        body = (
            '<div style="font-family:Arial,sans-serif;background:#f4f6fb;padding:30px;max-width:560px;margin:auto;">'
            '<div style="background:#ffffff;border-radius:10px;padding:28px;border:1px solid #dde3f5;">'
            '<h2 style="color:#1a6fef;margin-bottom:4px;">Meeting Invitation</h2>'
            '<p style="color:#6b7fa8;font-size:13px;margin-top:0;">You have been invited to an online meeting.</p>'
            '<hr style="border:none;border-top:1px solid #dde3f5;margin:20px 0;">'
            '<table style="width:100%;font-size:13px;color:#1a2340;">'
            '<tr><td style="padding:5px 0;color:#6b7fa8;width:120px;">Host</td><td><strong>' + str(employee_name) + '</strong></td></tr>'
            '<tr><td style="padding:5px 0;color:#6b7fa8;">Client</td><td><strong>' + str(client_name) + '</strong></td></tr>'
            '<tr><td style="padding:5px 0;color:#6b7fa8;">Scheduled</td><td>' + str(scheduled_time) + '</td></tr>'
            '<tr><td style="padding:5px 0;color:#6b7fa8;">Agenda</td><td>' + str(agenda) + '</td></tr>'
            '</table>'
            '<hr style="border:none;border-top:1px solid #dde3f5;margin:20px 0;">'
            '<p style="font-size:13px;color:#1a2340;margin-bottom:14px;">Click below to join:</p>'
            '<a href="' + str(join_url) + '" style="display:inline-block;background:#1a6fef;color:#ffffff;padding:13px 30px;border-radius:7px;text-decoration:none;font-weight:600;font-size:14px;">Join Meeting</a>'
            '<p style="font-size:12px;color:#6b7fa8;margin-top:14px;">Or copy this link:<br>'
            '<a href="' + str(join_url) + '" style="color:#1a6fef;word-break:break-all;">' + str(join_url) + '</a></p>'
            '<div style="background:#f4f6fb;border-radius:7px;padding:16px 18px;margin-top:18px;border:1px solid #dde3f5;">'
            '<div style="font-size:12px;color:#6b7fa8;margin-bottom:6px;">Meeting Password</div>'
            '<div style="font-size:28px;font-weight:700;letter-spacing:6px;color:#1a6fef;">' + str(password) + '</div>'
            '<div style="font-size:11px;color:#6b7fa8;margin-top:6px;">Enter this password when joining the meeting room.</div>'
            '</div>'
            '<hr style="border:none;border-top:1px solid #dde3f5;margin:20px 0;">'
            '<p style="font-size:11px;color:#6b7fa8;">Sent by XAI-ITD-DLP. Do not share your password.</p>'
            '</div></div>'
        )
        msg.attach(MIMEText(body, "html"))

        # Try STARTTLS (port 587) — same as OTP email
        try:
            with smtplib.SMTP(SMTP_SERVER, SMTP_PORT, timeout=15) as server:
                server.ehlo()
                server.starttls()
                server.ehlo()
                server.login(SMTP_EMAIL, SMTP_PASSWORD)
                server.sendmail(SMTP_EMAIL, to_email, msg.as_string())
        except Exception as e1:
            print("[EMAIL] STARTTLS failed (" + str(e1) + "), trying SSL port 465...")
            with smtplib.SMTP_SSL(SMTP_SERVER, 465, timeout=15) as server:
                server.login(SMTP_EMAIL, SMTP_PASSWORD)
                server.sendmail(SMTP_EMAIL, to_email, msg.as_string())

        print("[EMAIL] Meeting invite sent successfully to " + to_email)
        return True
    except Exception as e:
        print("[EMAIL ERROR] Meeting invite failed: " + str(e))
        return False


def login_required(roles=None):
    """
    Auth decorator that accepts EITHER:
      1. Flask session (normal browser login)
      2. X-Auth-Token header or ?token= query param (for multi-tab support)
    """
    def decorator(f):
        @wraps(f)
        def decorated(*args, **kwargs):
            token = request.headers.get("X-Auth-Token") or request.args.get("token")
            if token and token in active_tokens:
                email = active_tokens[token]
                user  = get_user_by_email(email)
                if user:
                    if roles and user["role"] not in roles:
                        return jsonify({"error": "Unauthorized"}), 403
                    request.auth_email = email
                    request.auth_role  = user["role"]
                    request.auth_name  = user["name"]
                    return f(*args, **kwargs)

            if "user_email" not in session:
                if request.path.startswith("/api/"):
                    return jsonify({"error": "Unauthorized"}), 403
                return redirect(url_for("login_page"))
            if roles and session.get("role") not in roles:
                return jsonify({"error": "Unauthorized"}), 403
            request.auth_email = session["user_email"]
            request.auth_role  = session.get("role")
            request.auth_name  = session.get("name", "")
            return f(*args, **kwargs)
        return decorated
    return decorator


def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        data = request.json or {}
        token = data.get("token")
        if not token or token not in active_tokens:
            return jsonify({"error": "Invalid token"}), 401
        return f(*args, **kwargs)
    return decorated


# --- Pages --------------------------------------------------------------------

@app.route("/")
def index():
    if "user_email" in session:
        role = session.get("role")
        if role == "admin":
            return redirect(url_for("admin_dashboard"))
        if role == "manager":
            return redirect(url_for("manager_dashboard"))
        return redirect(url_for("employee_dashboard"))
    return redirect(url_for("login_page"))


@app.route("/login")
def login_page():
    return render_template("login.html")


@app.route("/employee")
@login_required(roles=["employee"])
def employee_dashboard():
    user = get_user_by_email(session["user_email"])
    return render_template("employee.html", user=user, session_email=session["user_email"], session_role=session["role"], session_token=session.get("token",""))


@app.route("/manager")
@login_required(roles=["manager"])
def manager_dashboard():
    user = get_user_by_email(session["user_email"])
    return render_template("manager.html", user=user, session_email=session["user_email"], session_role=session["role"], session_token=session.get("token",""))


@app.route("/logout")
def logout():
    token = session.get("token")
    if token and token in active_tokens:
        del active_tokens[token]
    email = session.get("user_email")
    if email:
        log_activity(email, "LOGOUT", "User logged out.", session.get("location_str", "Unknown"), "-", "LOW")
    session.clear()
    return redirect(url_for("login_page"))


# --- Auth API -----------------------------------------------------------------

@app.route("/api/auth/request-otp", methods=["POST"])
def request_otp():
    data = request.json
    email = data.get("email", "").strip().lower()
    password = data.get("password", "")

    if not email or not password:
        return jsonify({"error": "Email and password required."}), 400

    selected_role = data.get("role", "").strip().lower()

    user = get_user_by_email(email)
    if not user or not verify_password(password, user["password"]):
        return jsonify({"error": "Invalid credentials."}), 401

    if not user.get("is_active", True):
        return jsonify({"error": "Account is deactivated. Contact administrator."}), 403

    if selected_role and user.get("role", "") != selected_role:
        return jsonify({"error": "Role mismatch. Your account is not registered as " + selected_role + "."}), 403

    ip = request.remote_addr
    location = get_location_from_ip(ip)
    city    = location.get("city", "Unknown")
    region  = location.get("region", "")
    country = location.get("country", "Unknown")

    role         = user.get("role", "employee")
    allowed_zones = user.get("allowed_locations", [])

    if role in ("manager", "admin"):
        loc_allowed = True
    elif not allowed_zones:
        loc_allowed = True
    elif city in ("Localhost", "Unknown", ""):
        # LOCAL MODE: geo-lookup failed or returned Unknown (common on localhost).
        # Allow login so local development is not blocked.
        # In production (app_deployed.py) this is set to False.
        loc_allowed = True
    else:
        candidates  = [s.lower() for s in [city, region, country] if s and s not in ("Unknown", "Local", "")]
        loc_allowed = False
        for zone in allowed_zones:
            zone_l = zone.lower()
            if any(zone_l in c or c in zone_l for c in candidates):
                loc_allowed = True
                break

    if not loc_allowed:
        detected = city + (" / " + region if region and region != city else "")
        loc_str  = city + ", " + country
        _locs = user.get("allowed_locations") or []
        approved_zones = ", ".join([z.get("city", "") if isinstance(z, dict) else str(z) for z in _locs]) or "None configured"
        xai_loc_reason = (
            "Your login was blocked because your current location ("
            + detected + ") is not in your approved location list. "
            "Your approved zones are: " + approved_zones + ". "
            "Contact your manager to add this location, or request travel mode if you are travelling."
        )
        log_security_event(
            email, "LOCATION_BLOCKED",
            "Login attempt from unauthorized location: " + loc_str,
            loc_str, ip, blocked=True
        )
        socketio.emit("new_event", {
            "type": "LOCATION_BLOCKED",
            "user": user["name"],
            "email": email,
            "detail": "Login blocked - unauthorized location: " + detected,
            "location": loc_str,
            "ip": ip,
            "time": datetime.utcnow().strftime("%H:%M:%S"),
            "risk": "HIGH",
            "blocked": True
        })
        return jsonify({
            "error":      "Access denied. Location detected: '" + detected + "'. Not in your allowed zones. Contact your administrator.",
            "xai_reason": xai_loc_reason
        }), 403

    otp = str(random.randint(100000, 999999))
    expiry = datetime.utcnow() + timedelta(seconds=OTP_EXPIRY_SECONDS)
    save_otp(email, otp, expiry)

    print("[OTP] " + email + " -> " + otp)

    threading.Thread(target=send_otp_email, args=(email, otp, user["name"]), daemon=True).start()

    return jsonify({"message": "OTP sent to " + email + ". Location verified: " + city + ", " + country + "."})


@app.route("/api/auth/verify-otp", methods=["POST"])
def verify_otp():
    data = request.json
    email = data.get("email", "").strip().lower()
    otp = data.get("otp", "").strip()

    if not email or not otp:
        return jsonify({"error": "Email and OTP required."}), 400

    otp_record = get_otp(email)
    if not otp_record:
        return jsonify({"error": "No OTP found. Please request again."}), 400

    if datetime.utcnow() > otp_record["expiry"]:
        delete_otp(email)
        return jsonify({"error": "OTP expired. Please request again."}), 400

    if otp_record["otp"] != otp:
        return jsonify({"error": "Incorrect OTP."}), 400

    delete_otp(email)

    user = get_user_by_email(email)
    ip = request.remote_addr
    location = get_location_from_ip(ip)
    city    = location.get("city", "Unknown")
    country = location.get("country", "Unknown")
    loc_str = city + ", " + country
    now     = datetime.utcnow()
    token   = secrets.token_hex(32)

    active_tokens[token] = email

    session["user_email"] = email
    session["role"]       = user["role"]
    session["name"]       = user["name"]
    session["token"]      = token
    session["location"]   = location
    session["location_str"] = loc_str

    # -- DEVICE FINGERPRINT ----------------------------------------
    fingerprint = data.get("fingerprint", {})
    device_mismatch, mismatch_reasons = False, []
    if fingerprint:
        device_mismatch, mismatch_reasons = check_device_mismatch(email, fingerprint)
        stored = get_device_profile(email)
        fp_os      = fingerprint.get("os", "")
        fp_browser = fingerprint.get("browser", "")
        fp_tz      = fingerprint.get("timezone", "")
        if not stored:
            save_device_profile(email, fingerprint)
            log_activity(email, "DEVICE_ENROLLED",
                "Device profile auto-enrolled: " + fp_os + " / " + fp_browser + " / " + fp_tz,
                loc_str, ip, "LOW")
        elif device_mismatch:
            mismatch_str = "; ".join(mismatch_reasons)
            log_activity(email, "DEVICE_MISMATCH",
                "Device changed: " + mismatch_str,
                loc_str, ip, "HIGH")
            log_security_event(email, "DEVICE_MISMATCH",
                "Suspicious device change detected: " + mismatch_str,
                loc_str, ip, blocked=False)
            emit_to_managers("new_event", {
                "type": "DEVICE_MISMATCH", "email": email, "user": user["name"],
                "detail": "New device login: " + mismatch_str,
                "location": loc_str, "risk": "HIGH", "blocked": False,
                "time": now.strftime("%H:%M:%S")
            })

    # -- IMPOSSIBLE TRAVEL -----------------------------------------
    travel_impossible, travel_detail = False, ""
    coords = {"lat": location.get("lat", 0), "lon": location.get("lon", 0)}
    is_travelling, travel_dest = is_travel_mode_active(email)

    if not is_travelling:
        travel_impossible, travel_detail = check_impossible_travel(email, coords, now)
        if travel_impossible:
            log_activity(email, "IMPOSSIBLE_TRAVEL",
                "Impossible travel detected: " + travel_detail,
                loc_str, ip, "HIGH")
            log_security_event(email, "IMPOSSIBLE_TRAVEL",
                travel_detail, loc_str, ip, blocked=False)
            emit_to_managers("new_event", {
                "type": "IMPOSSIBLE_TRAVEL", "email": email, "user": user["name"],
                "detail": "Impossible travel: " + travel_detail,
                "location": loc_str, "risk": "HIGH", "blocked": False,
                "time": now.strftime("%H:%M:%S")
            })
    else:
        log_activity(email, "TRAVEL_MODE_LOGIN",
            "Login during approved travel to " + travel_dest + " from " + loc_str,
            loc_str, ip, "LOW")

    update_user_login(email, loc_str, ip)
    update_login_coords(email, coords)

    login_risk = "HIGH" if (device_mismatch or travel_impossible) else "LOW"

    # Build login detail string without backslashes inside f-strings
    _fp_os = fingerprint.get("os", "") if fingerprint else ""
    _detail_parts = "Login from " + loc_str + " | IP: " + ip
    if device_mismatch:
        _detail_parts += " | ! DEVICE MISMATCH"
    if travel_impossible:
        _detail_parts += " | ! IMPOSSIBLE TRAVEL"
    if not device_mismatch and not travel_impossible:
        _detail_parts += " | Device: " + _fp_os

    log_activity(email, "LOGIN_SUCCESS", _detail_parts, loc_str, ip, risk_level=login_risk)
    if ai_score_now: ai_score_now(email, "LOGIN_SUCCESS", blocked=False, socketio_instance=socketio)

    _emit_detail = "Login from " + loc_str
    if device_mismatch:
        _emit_detail += " | ! NEW DEVICE"
    if travel_impossible:
        _emit_detail += " | ! IMPOSSIBLE TRAVEL"
    if is_travelling:
        _emit_detail += " | > TRAVEL MODE"

    emit_to_managers("new_event", {
        "type": "LOGIN_SUCCESS", "user": user["name"], "email": email,
        "role": user["role"],
        "detail": _emit_detail,
        "location": loc_str, "ip": ip,
        "time": now.strftime("%H:%M:%S"),
        "risk": login_risk, "blocked": False
    })

    print("[SESSION TOKEN] " + email + " -> " + token)

    if user["role"] in ("employee", "manager"):
        import json as _json
        token_path = os.path.join(os.path.dirname(__file__), "session_token.json")
        with open(token_path, "w") as tf:
            _json.dump({"email": email, "token": token, "name": user["name"], "role": user["role"]}, tf)
        print("[TOKEN FILE] Written to session_token.json for " + user["role"] + ": " + email)

        # Auto-launch agent from agent/start_agent.py if not already running
        try:
            root_dir = os.path.dirname(__file__)
            launcher = os.path.join(root_dir, "agent", "start_agent.py")
            if os.path.exists(launcher):
                already_running = False
                for proc in __import__("psutil").process_iter(["pid", "cmdline"]):
                    try:
                        cmd = " ".join(proc.info["cmdline"] or [])
                        if "start_agent" in cmd:
                            already_running = True
                            break
                    except Exception:
                        pass
                if not already_running:
                    _flags = subprocess.CREATE_NO_WINDOW if sys.platform == "win32" else 0
                    subprocess.Popen(
                        [sys.executable, launcher],
                        cwd=os.path.join(root_dir, "agent"),
                        stdout=subprocess.DEVNULL,
                        stderr=subprocess.DEVNULL,
                        creationflags=_flags
                    )
                    print("[AGENT] Auto-launched agent/start_agent.py for " + email)
                else:
                    print("[AGENT] Agent already running — new token written, will reload")
            else:
                print("[AGENT] agent/start_agent.py not found at: " + launcher)
        except Exception as _ae:
            print("[AGENT] Auto-launch error: " + str(_ae))

    _role = user["role"]
    if _role == "admin":
        _redirect = "/admin"
    elif _role == "manager":
        _redirect = "/manager"
    else:
        _redirect = "/employee"

    _xai_explanation = []
    _xai_score = 0
    _xai_level = "LOW"
    if device_mismatch:
        _xai_score += 40
        _xai_explanation.append("Login from an unrecognised device. Your browser, OS, or timezone changed from your enrolled profile. This may indicate account compromise or use of a new personal device.")
    if travel_impossible:
        _xai_score += 50
        _xai_explanation.append("Two logins were recorded from locations that are physically impossible to travel between in the elapsed time. This strongly suggests credential sharing or account compromise.")
    if _xai_score >= 70:
        _xai_level = "HIGH"
    elif _xai_score >= 35:
        _xai_level = "MEDIUM"

    return jsonify({
        "message":      "Login successful.",
        "role":         _role,
        "name":         user["name"],
        "token":        token,
        "location":     location,
        "location_str": loc_str,
        "redirect":     _redirect,
        "xai_login": {
            "risk_score":  _xai_score,
            "risk_level":  _xai_level,
            "explanation": _xai_explanation
        }
    })


# --- Agent API ----------------------------------------------------------------

@app.route("/api/agent/event", methods=["POST"])
@token_required
def receive_agent_event():
    data = request.json
    email = data.get("email")
    event_type = data.get("event_type")
    detail = data.get("detail")
    risk_level = data.get("risk_level", "LOW")
    blocked = data.get("blocked", False)

    user = get_user_by_email(email)
    if not user:
        return jsonify({"error": "User not found"}), 404

    ip = request.remote_addr
    location = get_location_from_ip(ip)
    loc_str = location.get("city", "") + ", " + location.get("country", "")

    if blocked or event_type in ("USB_INSERTED", "SCREENSHOT_BLOCKED"):
        log_security_event(email, event_type, detail, loc_str, ip, blocked=True)
    else:
        log_activity(email, event_type, detail, loc_str, ip, risk_level)

    # ── Real-time AI scoring ──────────────────────────────────────────────────
    if ai_score_now: ai_score_now(email, event_type, blocked=blocked, socketio_instance=socketio)

    socketio.emit("new_event", {
        "type": event_type,
        "user": user.get("name", email),
        "email": email,
        "detail": detail,
        "location": loc_str,
        "ip": ip,
        "time": datetime.utcnow().strftime("%H:%M:%S"),
        "risk": risk_level,
        "blocked": blocked
    })

    return jsonify({"status": "logged"})


# --- Manager APIs -------------------------------------------------------------

@app.route("/api/manager/logs")
@login_required(roles=["manager"])
def manager_logs():
    logs = get_all_logs(100)
    for log in logs:
        if isinstance(log.get("timestamp"), datetime):
            log["timestamp"] = log["timestamp"].strftime("%Y-%m-%d %H:%M:%S")
    return jsonify(logs)


@app.route("/api/manager/security-events")
@login_required(roles=["manager"])
def manager_security_events():
    events = get_all_security_events(100)
    for e in events:
        if isinstance(e.get("timestamp"), datetime):
            e["timestamp"] = e["timestamp"].strftime("%Y-%m-%d %H:%M:%S")
    return jsonify(events)


@app.route("/api/manager/users")
@login_required(roles=["manager"])
def manager_users():
    employees = get_all_employees()
    for u in employees:
        if isinstance(u.get("created_at"), datetime):
            u["created_at"] = u["created_at"].strftime("%Y-%m-%d")
        if isinstance(u.get("last_login_time"), datetime):
            u["last_login_time"] = u["last_login_time"].strftime("%Y-%m-%d %H:%M:%S")
    return jsonify(employees)


@app.route("/api/manager/stats")
@login_required(roles=["manager"])
def manager_stats():
    from models.user import logs_col, events_col
    total_events = logs_col.count_documents({})
    high_risk = logs_col.count_documents({"risk_level": "HIGH"})
    blocked = events_col.count_documents({"blocked": True})
    active_users = len(active_tokens)
    location_blocks = events_col.count_documents({"action": "LOCATION_BLOCKED"})
    return jsonify({
        "total_events": total_events,
        "high_risk": high_risk,
        "blocked_actions": blocked,
        "active_users": active_users,
        "location_blocks": location_blocks
    })


@app.route("/api/manager/all-logs")
@login_required(roles=["manager"])
def manager_all_logs():
    from models.user import logs_col
    docs = list(logs_col.find({}, {"_id": 0}).sort("timestamp", -1).limit(500))
    for d in docs:
        if hasattr(d.get("timestamp"), "strftime"):
            d["timestamp"] = d["timestamp"].strftime("%Y-%m-%d %H:%M:%S")
    return jsonify(docs)


@app.route("/api/manager/login-history/<email>")
@login_required(roles=["manager"])
def employee_login_history(email):
    history = get_login_history(email)
    return jsonify(history)


@app.route("/api/manager/set-locations", methods=["POST"])
@login_required(roles=["manager"])
def manager_set_locations():
    data = request.json
    email = data.get("email")
    locations = data.get("locations", [])
    if not email:
        return jsonify({"error": "Email required"}), 400
    set_allowed_locations(email, locations)
    log_activity(
        request.auth_email, "LOCATION_POLICY_UPDATED",
        "Updated allowed locations for " + email + ": " + str(locations),
        "Manager Dashboard", "internal", "LOW"
    )
    return jsonify({"message": "Allowed locations updated for " + email + "."})


# --- Employee APIs ------------------------------------------------------------

@app.route("/api/employee/my-logs")
@login_required(roles=["employee"])
def employee_logs():
    logs = get_user_logs(request.auth_email, 30)
    for log in logs:
        if isinstance(log.get("timestamp"), datetime):
            log["timestamp"] = log["timestamp"].strftime("%Y-%m-%d %H:%M:%S")
    return jsonify(logs)


@app.route("/api/employee/context")
@login_required(roles=["employee"])
def employee_context():
    user = get_user_by_email(request.auth_email)
    loc = session.get("location") if "user_email" in session and session.get("user_email") == request.auth_email else None
    loc_str = session.get("location_str", "") if "user_email" in session and session.get("user_email") == request.auth_email else ""
    return jsonify({
        "name": request.auth_name,
        "email": request.auth_email,
        "role": request.auth_role,
        "location": loc or ({"city": user.get("last_login_location",""), "country":""} if user else {}),
        "location_str": loc_str or (user.get("last_login_location","") if user else "")
    })


# ===============================================================
# ADMIN - PAGES & API
# ===============================================================

@app.route("/admin")
def admin_dashboard():
    if "user_email" not in session or session.get("role") != "admin":
        return redirect(url_for("login_page"))
    user = get_user_by_email(session["user_email"])
    return render_template("admin.html", user=user, session_token=session.get("token",""))


@app.route("/api/admin/stats")
@login_required(roles=["admin"])
def admin_stats():
    stats = get_system_stats()
    stats["active_sessions"] = len(active_tokens)
    return jsonify(stats)


@app.route("/api/admin/all-logs")
@login_required(roles=["admin"])
def admin_all_logs():
    return jsonify(get_all_logs_unrestricted(200))


@app.route("/api/admin/all-security-events")
@login_required(roles=["admin"])
def admin_all_security_events():
    return jsonify(get_all_security_events_unrestricted(200))


@app.route("/api/admin/managers")
@login_required(roles=["admin"])
def admin_list_managers():
    return jsonify(get_all_users_by_role("manager"))


@app.route("/api/admin/employees")
@login_required(roles=["admin"])
def admin_list_employees():
    return jsonify(get_all_users_by_role("employee"))


@app.route("/api/admin/add-user", methods=["POST"])
@login_required(roles=["admin"])
def admin_add_user():
    data       = request.json
    name       = data.get("name", "").strip()
    email      = data.get("email", "").strip().lower()
    password   = data.get("password", "").strip()
    role       = data.get("role", "employee")
    department = data.get("department", "IT").strip()

    if not name or not email or not password:
        return jsonify({"error": "Name, email and password are required."}), 400
    if role not in ("employee", "manager"):
        return jsonify({"error": "Role must be employee or manager."}), 400
    if get_user_by_email(email):
        return jsonify({"error": "User with email " + email + " already exists."}), 409

    create_user(name, email, password, role, department)

    log_activity(request.auth_email, "ADMIN_USER_CREATED",
                 "Created " + role + " account: " + name + " (" + email + ")",
                 "Admin Dashboard", "internal", "LOW")

    return jsonify({"message": role.capitalize() + " '" + name + "' created successfully."})


@app.route("/api/admin/deactivate-user", methods=["POST"])
@login_required(roles=["admin"])
def admin_deactivate_user():
    email = request.json.get("email")
    if not email:
        return jsonify({"error": "Email required."}), 400
    deactivate_user(email)
    log_activity(request.auth_email, "ADMIN_USER_DEACTIVATED",
                 "Deactivated account: " + email,
                 "Admin Dashboard", "internal", "MEDIUM")
    return jsonify({"message": "Account " + email + " deactivated."})


@app.route("/api/admin/reactivate-user", methods=["POST"])
@login_required(roles=["admin"])
def admin_reactivate_user():
    email = request.json.get("email")
    if not email:
        return jsonify({"error": "Email required."}), 400
    reactivate_user(email)
    log_activity(request.auth_email, "ADMIN_USER_REACTIVATED",
                 "Reactivated account: " + email,
                 "Admin Dashboard", "internal", "LOW")
    return jsonify({"message": "Account " + email + " reactivated."})


@app.route("/api/admin/user-logs/<email>")
@login_required(roles=["admin"])
def admin_user_logs(email):
    logs   = get_all_logs_for_user(email, 100)
    events = get_all_security_events_for_user(email, 50)
    user   = get_user_by_email(email)
    history = []
    if user:
        history = user.get("login_history", [])
        for h in history:
            if hasattr(h.get("time"), "strftime"):
                h["time"] = h["time"].strftime("%Y-%m-%d %H:%M:%S")
        history = list(reversed(history))
    return jsonify({"logs": logs, "security_events": events, "login_history": history})


@app.route("/api/admin/travel-tracking")
@login_required(roles=["admin"])
def admin_travel_tracking():
    reqs = get_all_travel_requests()
    for r in reqs:
        tm = r.get("travel_mode", {}) or {}
        for k in ("requested_at", "approved_at", "start_date", "end_date"):
            if hasattr(tm.get(k), "strftime"):
                tm[k] = tm[k].strftime("%Y-%m-%d %H:%M")
    return jsonify(reqs)


# ===============================================================
# FILE SHARING - MANAGER UPLOAD
# ===============================================================

@app.route("/api/manager/upload-file", methods=["POST"])
@login_required(roles=["manager"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file provided."}), 400
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected."}), 400
    if not allowed_file(file.filename):
        return jsonify({"error": "File type not allowed."}), 400

    visibility     = request.form.get("visibility", "public")
    allowed_emails = request.form.getlist("allowed_emails")

    ext           = file.filename.rsplit(".", 1)[1].lower()
    unique_name   = uuid.uuid4().hex + "." + ext
    file_bytes    = file.read()
    file_size     = len(file_bytes)
    save_file_to_gridfs(unique_name, file_bytes, file.content_type)

    record = save_file_record(
        filename      = unique_name,
        original_name = secure_filename(file.filename),
        file_size     = file_size,
        file_type     = file.content_type,
        visibility    = visibility,
        allowed_emails= allowed_emails,
        uploaded_by   = request.auth_email
    )

    emit_to_employees("file_uploaded", {
        "file_id":    record["_id"],
        "name":       record["original_name"],
        "visibility": visibility,
        "time":       datetime.utcnow().strftime("%H:%M:%S")
    })
    emit_to_managers("new_event", {
        "type":     "FILE_UPLOADED",
        "email":    request.auth_email,
        "user":     request.auth_email,
        "detail":   "Uploaded file: " + record["original_name"] + " (" + visibility + ")",
        "location": "Manager Dashboard",
        "risk":     "LOW",
        "blocked":  False,
        "time":     datetime.utcnow().strftime("%H:%M:%S")
    })

    log_activity(request.auth_email, "FILE_UPLOADED",
                 "Uploaded: " + record["original_name"] + " (" + visibility + ")",
                 "Manager Dashboard", "internal", "LOW")

    return jsonify({"message": "File uploaded successfully.", "file": record})


@app.route("/api/manager/files")
@login_required(roles=["manager"])
def manager_list_files():
    return jsonify(get_all_files())


@app.route("/api/manager/view-file/<file_id>", methods=["GET", "HEAD"])
@login_required(roles=["manager"])
def manager_view_file(file_id):
    rec = get_file_by_id_unrestricted(file_id)
    if not rec:
        return jsonify({"error": "File not found."}), 404
    mime  = rec.get("file_type") or "application/octet-stream"
    fname = rec.get("original_name", "file")
    log_activity(request.auth_email, "FILE_VIEWED",
                 "Manager viewed: " + fname, "Manager Dashboard",
                 request.remote_addr, "LOW")
    from flask import make_response, Response
    if request.method == "HEAD":
        resp = make_response("", 200)
        resp.headers["Content-Type"] = mime
        return resp
    file_bytes = get_file_from_gridfs(rec["filename"])
    if file_bytes is None:
        return jsonify({"error": "File not found in storage."}), 404
    resp = make_response(Response(file_bytes, mimetype=mime))
    resp.headers["Content-Disposition"] = 'inline; filename="' + fname + '"'
    resp.headers["Cache-Control"] = "no-store, no-cache, must-revalidate"
    return resp


@app.route("/api/manager/recent-files")
@login_required(roles=["manager"])
def manager_recent_files():
    from models.files import shared_files_col, approval_col
    from bson import ObjectId
    email   = request.auth_email
    results = []
    seen    = set()
    for f in shared_files_col.find({"uploaded_by": email, "is_active": True}).sort("uploaded_at", -1).limit(50):
        fid = str(f["_id"])
        if fid in seen:
            continue
        seen.add(fid)
        f["_id"]       = fid
        f["direction"] = "sent"
        if hasattr(f.get("uploaded_at"), "strftime"):
            f["uploaded_at"] = f["uploaded_at"].strftime("%Y-%m-%d %H:%M")
        results.append(f)
    for req in approval_col.find({"forward_to": email, "request_type": "forward", "status": "approved"}).sort("resolved_at", -1).limit(30):
        fid = str(req.get("file_id", ""))
        if fid in seen:
            continue
        try:
            fobj = shared_files_col.find_one({"_id": ObjectId(fid), "is_active": True})
            if fobj:
                fobj["_id"]       = str(fobj["_id"])
                fobj["direction"] = "received"
                fobj["forwarded"] = True
                if hasattr(fobj.get("uploaded_at"), "strftime"):
                    fobj["uploaded_at"] = fobj["uploaded_at"].strftime("%Y-%m-%d %H:%M")
                results.append(fobj)
                seen.add(fid)
        except Exception:
            pass
    return jsonify(results)


@app.route("/api/manager/delete-file/<file_id>", methods=["DELETE"])
@login_required(roles=["manager"])
def manager_delete_file(file_id):
    rec = get_file_by_id(file_id)
    if not rec:
        return jsonify({"error": "File not found."}), 404
    delete_file_from_gridfs(rec["filename"])
    delete_file_record(file_id)
    return jsonify({"message": "File deleted."})


@app.route("/api/manager/update-file-visibility", methods=["POST"])
@login_required(roles=["manager"])
def manager_update_file_visibility():
    from models.files import shared_files_col
    from bson import ObjectId
    data           = request.json or {}
    file_id        = data.get("file_id", "").strip()
    visibility     = data.get("visibility", "public")
    allowed_emails = data.get("allowed_emails", [])

    if not file_id:
        return jsonify({"error": "file_id is required."}), 400
    if visibility not in ("public", "private"):
        return jsonify({"error": "visibility must be 'public' or 'private'."}), 400

    try:
        result = shared_files_col.update_one(
            {"_id": ObjectId(file_id), "is_active": True},
            {"$set": {
                "visibility":     visibility,
                "allowed_emails": allowed_emails if visibility == "private" else []
            }}
        )
        if result.matched_count == 0:
            return jsonify({"error": "File not found."}), 404

        recipients = ", ".join(allowed_emails) if allowed_emails else "all employees"
        log_activity(
            request.auth_email, "FILE_SEND",
            "Sent file " + file_id + " → " + visibility + " (" + recipients + ")",
            "Manager Dashboard", "internal", "LOW"
        )
        emit_to_employees("file_updated", {
            "file_id":    file_id,
            "visibility": visibility,
            "time":       datetime.utcnow().strftime("%H:%M:%S")
        })
        return jsonify({"message": "File sent successfully.", "visibility": visibility})
    except Exception as e:
        return jsonify({"error": "Server error: " + str(e)}), 500


@app.route("/api/manager/pending-approvals")
@login_required(roles=["manager"])
def manager_pending_approvals():
    return jsonify(get_pending_approvals())


@app.route("/api/manager/all-approvals")
@login_required(roles=["manager"])
def manager_all_approvals():
    return jsonify(get_all_approvals())


@app.route("/api/manager/resolve-approval", methods=["POST"])
@login_required(roles=["manager"])
def manager_resolve_approval():
    data          = request.json or {}
    request_id    = data.get("request_id")
    status        = data.get("status")
    reject_reason = data.get("reason", "")

    if status not in ("approved", "rejected"):
        return jsonify({"error": "Invalid status."}), 400

    req = resolve_approval(request_id, status, request.auth_email, reject_reason)
    if not req:
        return jsonify({"error": "Request not found."}), 404

    requested_by = req.get("requested_by", "")
    request_type = req.get("request_type", "")
    file_id      = req.get("file_id", "")
    forward_to   = req.get("forward_to", "")

    if status == "approved" and file_id:
        from models.files import shared_files_col
        from bson import ObjectId
        try:
            if request_type == "download":
                shared_files_col.update_one(
                    {"_id": ObjectId(file_id)},
                    {"$addToSet": {"allowed_emails": requested_by}}
                )
            elif request_type == "forward" and forward_to:
                shared_files_col.update_one(
                    {"_id": ObjectId(file_id)},
                    {"$addToSet": {"allowed_emails": forward_to}}
                )
                rec = get_file_by_id(file_id)
                if rec:
                    emit_to_user(forward_to, "file_uploaded", {
                        "file_id":    file_id,
                        "name":       rec.get("original_name", ""),
                        "sent_by":    requested_by,
                        "visibility": "private",
                        "time":       datetime.utcnow().strftime("%H:%M:%S")
                    })
        except Exception as e:
            app.logger.warning("resolve_approval: failed to update allowed_emails: " + str(e))

    emit_to_user(requested_by, "approval_resolved", {
        "request_id":   request_id,
        "file_name":    req.get("file_name"),
        "request_type": request_type,
        "status":       status,
        "reason":       reject_reason,
        "email":        requested_by
    })
    emit_to_managers("new_event", {
        "type":     "APPROVAL_" + status.upper(),
        "email":    request.auth_email,
        "user":     request.auth_email,
        "detail":   status.capitalize() + " " + request_type + " for " + str(req.get("file_name")) + " by " + requested_by,
        "location": "Manager Dashboard",
        "risk":     "LOW",
        "blocked":  False,
        "time":     datetime.utcnow().strftime("%H:%M:%S")
    })

    log_activity(request.auth_email, "APPROVAL_" + status.upper(),
                 status.capitalize() + " " + request_type + " for " + str(req.get("file_name")) + " by " + str(req.get("requested_by")),
                 "Manager Dashboard", "internal", "LOW")

    return jsonify({"message": "Request " + status + ".", "request": req})


# ===============================================================
# FILE SHARING - EMPLOYEE ACCESS
# ===============================================================

@app.route("/api/employee/files")
@login_required(roles=["employee"])
def employee_list_files():
    files = get_files_for_employee(request.auth_email)
    return jsonify(files)


@app.route("/api/employee/view-file/<file_id>", methods=["GET", "HEAD"])
@login_required(roles=["employee"])
def employee_view_file(file_id):
    rec = get_file_by_id(file_id)
    if not rec:
        return jsonify({"error": "File not found."}), 404

    email = request.auth_email

    allowed = False
    if rec["visibility"] == "public":
        allowed = True
    elif email in rec.get("allowed_emails", []):
        allowed = True
    else:
        from models.files import approval_col
        fwd = approval_col.find_one({
            "file_id": str(file_id),
            "forward_to": email,
            "request_type": "forward",
            "status": "approved"
        })
        if fwd:
            allowed = True

    if not allowed:
        return jsonify({"error": "Access denied."}), 403

    log_activity(email, "FILE_VIEWED",
                 "Viewed file: " + rec["original_name"],
                 "Employee Dashboard", request.remote_addr, "LOW")

    mime  = rec.get("file_type") or "application/octet-stream"
    fname = rec.get("original_name", "file")

    from flask import make_response, Response
    if request.method == "HEAD":
        resp = make_response("", 200)
        resp.headers["Content-Type"] = mime
        return resp
    file_bytes = get_file_from_gridfs(rec["filename"])
    if file_bytes is None:
        return jsonify({"error": "File not found in storage."}), 404
    resp = make_response(Response(file_bytes, mimetype=mime))
    resp.headers["Content-Disposition"] = 'inline; filename="' + fname + '"'
    resp.headers["Cache-Control"] = "no-store, no-cache, must-revalidate"
    return resp


# ──────────────────────────────────────────────────────────────
# SHARED FILE PREVIEW HELPER
# ──────────────────────────────────────────────────────────────
def _generate_file_preview_html(file_bytes, original_name, mime):
    import html as _html
    import io as _io
    ext = os.path.splitext(original_name)[1].lower().lstrip(".")

    text_exts = {"txt","log","md","py","js","ts","css","html","xml","yaml","yml",
                 "json","sh","bat","ini","cfg","env","csv","tsv"}
    if ext in text_exts or "text" in (mime or "") or "json" in (mime or "") or "xml" in (mime or ""):
        try:
            raw = file_bytes.decode("utf-8", errors="replace")[:80000]
            if ext == "csv" or "csv" in (mime or ""):
                import csv
                reader = csv.reader(_io.StringIO(raw))
                rows = list(reader)[:500]
                thead = "<tr>" + "".join("<th>" + _html.escape(str(c)) + "</th>" for c in (rows[0] if rows else [])) + "</tr>"
                tbody = "".join("<tr>" + "".join("<td>" + _html.escape(str(c)) + "</td>" for c in row) + "</tr>" for row in rows[1:100])
                html_out = (
                    "<style>"
                    "table{border-collapse:collapse;width:100%;font-family:Share Tech Mono,monospace;font-size:11px;}"
                    "th{background:#1a1a2e;color:#00d4ff;padding:6px 10px;border:1px solid #333;text-align:left;}"
                    "td{padding:5px 10px;border:1px solid #222;color:#ccc;}"
                    "tr:nth-child(even){background:#111;}"
                    "</style>"
                    "<div style='padding:12px;overflow:auto;max-height:520px;'>"
                    "<table><thead>" + thead + "</thead><tbody>" + tbody + "</tbody></table>"
                    "<div style='color:#555;font-size:10px;font-family:Share Tech Mono,monospace;margin-top:8px;'>Showing up to 100 rows</div>"
                    "</div>"
                )
            else:
                escaped = _html.escape(raw)
                html_out = (
                    "<style>pre{background:#0d0d1a;color:#b0c4de;padding:16px;margin:0;"
                    "font-family:Share Tech Mono,monospace;font-size:12px;overflow:auto;"
                    "max-height:520px;white-space:pre-wrap;word-break:break-all;}</style>"
                    "<pre>" + escaped + "</pre>"
                )
            return html_out, 200
        except Exception:
            return "", 415

    if ext in ("docx", "doc"):
        try:
            import docx as _docx
            doc = _docx.Document(_io.BytesIO(file_bytes))
            parts = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    parts.append("<br>")
                    continue
                style = (para.style.name or "").lower()
                esc = _html.escape(text)
                if "heading 1" in style:
                    parts.append("<h1>" + esc + "</h1>")
                elif "heading 2" in style:
                    parts.append("<h2>" + esc + "</h2>")
                elif "heading 3" in style:
                    parts.append("<h3>" + esc + "</h3>")
                else:
                    runs_html = ""
                    for run in para.runs:
                        r = _html.escape(run.text)
                        if run.bold and run.italic:
                            r = "<strong><em>" + r + "</em></strong>"
                        elif run.bold:
                            r = "<strong>" + r + "</strong>"
                        elif run.italic:
                            r = "<em>" + r + "</em>"
                        runs_html += r
                    parts.append("<p>" + (runs_html or esc) + "</p>")
            for table in doc.tables:
                rows_html = ""
                for i, row in enumerate(table.rows):
                    tag = "th" if i == 0 else "td"
                    cells = "".join(
                        "<" + tag + ">" + _html.escape(cell.text.strip()) + "</" + tag + ">"
                        for cell in row.cells
                    )
                    rows_html += "<tr>" + cells + "</tr>"
                parts.append("<table>" + rows_html + "</table>")
            styled = (
                "<style>"
                "body,div{font-family:Segoe UI,Arial,sans-serif;color:#d0d8e8;background:transparent;}"
                "h1,h2,h3{color:#60aaff;margin:12px 0 6px;}"
                "h1{font-size:18px;}h2{font-size:15px;}h3{font-size:13px;}"
                "table{border-collapse:collapse;width:100%;margin:10px 0;}"
                "td,th{border:1px solid #333;padding:6px 10px;color:#ccc;font-size:12px;}"
                "th{background:#1a1a2e;color:#00d4ff;}"
                "p{margin:4px 0;font-size:13px;line-height:1.6;}"
                "strong{color:#fff;}em{color:#aac4ff;}"
                "</style>"
                "<div style='padding:20px;overflow:auto;max-height:520px;'>" + "".join(parts) + "</div>"
            )
            return styled, 200
        except Exception:
            return "", 415

    if ext in ("xlsx", "xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(_io.BytesIO(file_bytes), read_only=True, data_only=True)
            sheets_html = ""
            for sheet_name in wb.sheetnames[:5]:
                ws = wb[sheet_name]
                rows_html = ""
                is_first = True
                for row in ws.iter_rows(max_row=200, values_only=True):
                    tag = "th" if is_first else "td"
                    cells = "".join(
                        "<" + tag + ">" + _html.escape(str(c) if c is not None else "") + "</" + tag + ">"
                        for c in row
                    )
                    rows_html += "<tr>" + cells + "</tr>"
                    is_first = False
                sheets_html += (
                    "<div style='margin-bottom:20px;'>"
                    "<div style='font-family:Share Tech Mono,monospace;font-size:10px;color:#00d4ff;"
                    "letter-spacing:2px;margin-bottom:8px;'>" + _html.escape(sheet_name) + "</div>"
                    "<div style='overflow-x:auto;'><table>" + rows_html + "</table></div>"
                    "<div style='color:#555;font-size:10px;font-family:Share Tech Mono,monospace;margin-top:4px;'>Showing up to 200 rows</div>"
                    "</div>"
                )
            html_out = (
                "<style>"
                "table{border-collapse:collapse;font-family:Share Tech Mono,monospace;font-size:11px;}"
                "th{background:#1a1a2e;color:#00d4ff;padding:6px 10px;border:1px solid #333;text-align:left;}"
                "td{padding:5px 10px;border:1px solid #222;color:#ccc;}"
                "tr:nth-child(even){background:#0a0a18;}"
                "</style>"
                "<div style='padding:14px;overflow:auto;max-height:520px;'>" + sheets_html + "</div>"
            )
            return html_out, 200
        except ImportError:
            pass
        except Exception:
            return "", 415

    if ext in ("pptx", "ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(_io.BytesIO(file_bytes))
            slides_html = ""
            for i, slide in enumerate(prs.slides[:30]):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                slide_text = "<br>".join(_html.escape(t) for t in texts)
                slides_html += (
                    "<div style='background:#0d0d1a;border:1px solid #1a1a2e;border-radius:8px;"
                    "padding:16px;margin-bottom:12px;'>"
                    "<div style='font-family:Share Tech Mono,monospace;font-size:9px;color:#555;"
                    "letter-spacing:2px;margin-bottom:8px;'>SLIDE " + str(i+1) + "</div>"
                    "<div style='color:#d0d8e8;font-size:13px;line-height:1.6;'>" + (slide_text or "<em style='color:#444;'>Empty slide</em>") + "</div>"
                    "</div>"
                )
            html_out = (
                "<div style='padding:14px;overflow:auto;max-height:520px;'>"
                + (slides_html or "<div style='color:#555;text-align:center;padding:40px;'>No text content found.</div>")
                + "</div>"
            )
            return html_out, 200
        except ImportError:
            pass
        except Exception:
            return "", 415

    return "", 415


@app.route("/api/employee/preview-file/<file_id>")
@login_required(roles=["employee"])
def employee_preview_file(file_id):
    rec = get_file_by_id(file_id)
    if not rec:
        return jsonify({"error": "File not found."}), 404

    email = request.auth_email
    allowed = rec["visibility"] == "public" or email in rec.get("allowed_emails", [])
    if not allowed:
        from models.files import approval_col
        fwd = approval_col.find_one({
            "file_id": str(file_id), "forward_to": email,
            "request_type": "forward", "status": "approved"
        })
        allowed = bool(fwd)
    if not allowed:
        return jsonify({"error": "Access denied."}), 403

    file_bytes = get_file_from_gridfs(rec["filename"])
    if file_bytes is None:
        return jsonify({"error": "File not found in storage."}), 404

    html_out, status = _generate_file_preview_html(file_bytes, rec.get("original_name", "file"), rec.get("file_type", ""))
    if status != 200:
        return "", 415
    from flask import Response
    return Response(html_out, status=200, mimetype="text/html")


@app.route("/api/manager/preview-file/<file_id>")
@login_required(roles=["manager"])
def manager_preview_file(file_id):
    rec = get_file_by_id_unrestricted(file_id)
    if not rec:
        return jsonify({"error": "File not found."}), 404

    file_bytes = get_file_from_gridfs(rec["filename"])
    if file_bytes is None:
        return jsonify({"error": "File not found in storage."}), 404

    html_out, status = _generate_file_preview_html(file_bytes, rec.get("original_name", "file"), rec.get("file_type", ""))
    if status != 200:
        return "", 415
    from flask import Response
    return Response(html_out, status=200, mimetype="text/html")


@app.route("/api/employee/request-access", methods=["POST"])
@login_required(roles=["employee"])
def employee_request_access():
    data         = request.json or {}
    file_id      = data.get("file_id")
    request_type = data.get("request_type")
    forward_to   = data.get("forward_to")
    file_name    = data.get("file_name", "")

    if not file_id or not request_type:
        return jsonify({"error": "file_id and request_type are required."}), 400

    rec = get_file_by_id(file_id)
    if not rec:
        return jsonify({"error": "File not found."}), 404

    email = request.auth_email

    if request_type == "download":
        if rec["visibility"] == "public" or email in rec.get("allowed_emails", []):
            return jsonify({"message": "You already have access. No approval needed.", "already_allowed": True}), 200

    req, is_new = create_approval_request(
        file_id      = file_id,
        file_name    = rec.get("original_name", file_name),
        requested_by = email,
        request_type = request_type,
        forward_to   = forward_to
    )

    if not is_new:
        return jsonify({"message": "Request already pending. Please wait for manager approval."}), 200

    emit_to_managers("new_approval_request", {
        "request_id":   req["_id"],
        "file_name":    rec.get("original_name", file_name),
        "requested_by": email,
        "request_type": request_type,
        "forward_to":   forward_to,
        "time":         datetime.utcnow().strftime("%H:%M:%S")
    })

    log_activity(email, "ACCESS_REQUESTED_" + request_type.upper(),
                 "Requested " + request_type + " for: " + rec.get("original_name", file_name),
                 "Employee Dashboard", request.remote_addr, "MEDIUM")

    return jsonify({"message": "Request sent to manager for approval.", "request_id": req["_id"]})


@app.route("/api/employee/download-file/<file_id>")
@login_required(roles=["employee"])
def employee_download_file(file_id):
    rec = get_file_by_id(file_id)
    if not rec:
        return jsonify({"error": "File not found."}), 404

    email = request.auth_email

    from models.files import approval_col

    approved = approval_col.find_one({
        "file_id":      file_id,
        "requested_by": email,
        "request_type": "download",
        "status":       "approved"
    })

    if not approved:
        return jsonify({"error": "Download not approved. Please request access first."}), 403

    log_activity(email, "FILE_DOWNLOADED",
                 "Downloaded: " + rec["original_name"],
                 "Employee Dashboard", request.remote_addr, "MEDIUM")

    file_bytes = get_file_from_gridfs(rec["filename"])
    if file_bytes is None:
        return jsonify({"error": "File not found in storage."}), 404
    from flask import make_response, Response
    resp = make_response(Response(file_bytes, mimetype=rec.get("file_type","application/octet-stream")))
    resp.headers["Content-Disposition"] = 'attachment; filename="' + rec["original_name"] + '"'
    return resp


@app.route("/api/employee/my-file-approvals")
@login_required(roles=["employee"])
def employee_my_file_approvals():
    from models.files import approval_col
    email = request.auth_email
    reqs = list(approval_col.find({"requested_by": email}))
    result = {}
    for r in reqs:
        fid = r.get("file_id", "")
        rtype = r.get("request_type", "")
        status = r.get("status", "pending")
        if fid not in result:
            result[fid] = {}
        existing = result[fid].get(rtype)
        if existing != "approved":
            result[fid][rtype] = status
            if rtype == "forward":
                result[fid]["forward_to"] = r.get("forward_to", "")
    return jsonify(result)


@app.route("/api/employee/other-employees")
@login_required(roles=["employee"])
def employee_other_employees():
    all_emps = get_all_employees()
    others = [e for e in all_emps if e.get("email") != request.auth_email]
    return jsonify([{"name": e["name"], "email": e["email"]} for e in others])


@app.route("/api/employee/my-requests")
@login_required(roles=["employee"])
def employee_my_requests():
    return jsonify(get_employee_requests(request.auth_email))


@app.route("/api/agent/event-browser", methods=["POST"])
@login_required(roles=["employee"])
def agent_event_browser():
    data       = request.json or {}
    event_type = data.get("event_type", "BROWSER_EVENT")
    detail     = data.get("detail", "")
    risk       = data.get("risk", "HIGH")
    blocked    = data.get("blocked", True)

    log_activity(request.auth_email, event_type, detail, "Browser", request.remote_addr, risk)
    log_security_event(request.auth_email, event_type, detail, "Browser", request.remote_addr, blocked)

    emit_to_managers("new_event", {
        "type":     event_type,
        "email":    request.auth_email,
        "user":     request.auth_email,
        "detail":   detail,
        "location": "Browser",
        "risk":     risk,
        "blocked":  blocked,
        "time":     datetime.utcnow().strftime("%H:%M:%S")
    })
    return jsonify({"status": "logged"})


@app.route("/api/employee/my-uploads")
@login_required(roles=["employee"])
def employee_my_uploads():
    from models.files import shared_files_col
    email = request.auth_email
    docs = list(shared_files_col.find(
        {"uploaded_by": email, "is_active": True},
        {"_id": 1, "original_name": 1, "file_size": 1, "uploaded_at": 1,
         "scan_clean": 1, "scan_engine": 1}
    ).sort("uploaded_at", -1).limit(50))
    for d in docs:
        d["_id"] = str(d["_id"])
        if hasattr(d.get("uploaded_at"), "strftime"):
            d["uploaded_at"] = d["uploaded_at"].strftime("%Y-%m-%d %H:%M")
    return jsonify(docs)


# ===============================================================
# TRAVEL MODE ROUTES
# ===============================================================

@app.route("/api/employee/request-travel", methods=["POST"])
@login_required(roles=["employee"])
def employee_request_travel():
    data        = request.json or {}
    source      = data.get("source", "").strip()
    destination = data.get("destination", "").strip()
    start_date  = data.get("start_date", "").strip()
    end_date    = data.get("end_date", "").strip()
    reason      = data.get("reason", "").strip()
    src_coords  = data.get("src_coords")
    dst_coords  = data.get("dst_coords")

    if not source:
        return jsonify({"error": "Source location required. Pin both points on the map."}), 400
    if not destination:
        return jsonify({"error": "Destination location required."}), 400
    if not start_date or not end_date:
        return jsonify({"error": "Start and end date/time required."}), 400
    if not reason:
        return jsonify({"error": "Reason for travel is required."}), 400
    try:
        start_dt = datetime.strptime(start_date, "%Y-%m-%dT%H:%M")
        end_dt   = datetime.strptime(end_date,   "%Y-%m-%dT%H:%M")
    except Exception:
        return jsonify({"error": "Invalid date format. Expected YYYY-MM-DDTHH:MM."}), 400
    if end_dt <= start_dt:
        return jsonify({"error": "End date/time must be after start."}), 400
    if (end_dt - start_dt).days > 30:
        return jsonify({"error": "Travel mode cannot exceed 30 days."}), 400

    if src_coords:
        try:
            src_coords = {"lat": float(src_coords["lat"]), "lon": float(src_coords["lon"])}
        except Exception:
            src_coords = None
    if dst_coords:
        try:
            dst_coords = {"lat": float(dst_coords["lat"]), "lon": float(dst_coords["lon"])}
        except Exception:
            dst_coords = None

    request_travel_mode(
        request.auth_email, destination, start_dt, end_dt, reason,
        source=source, src_coords=src_coords, dst_coords=dst_coords
    )
    log_activity(request.auth_email, "TRAVEL_MODE_REQUESTED",
        "Travel requested: " + source + " -> " + destination + " (" + start_date + " to " + end_date + "). Reason: " + reason,
        "Employee Dashboard", request.remote_addr, "MEDIUM")

    emit_to_managers("travel_request", {
        "type":        "travel_mode",
        "email":       request.auth_email,
        "source":      source,
        "destination": destination,
        "start_date":  start_date,
        "end_date":    end_date,
        "reason":      reason,
        "src_coords":  src_coords,
        "dst_coords":  dst_coords,
        "time":        datetime.utcnow().strftime("%H:%M:%S")
    })
    return jsonify({"message": "Travel mode request submitted. Awaiting manager approval."})


@app.route("/api/employee/travel-status")
@login_required(roles=["employee"])
def employee_travel_status():
    tm = get_travel_mode(request.auth_email)
    return jsonify(tm or {})


@app.route("/api/employee/travel-history")
@login_required(roles=["employee"])
def employee_travel_history():
    history = get_travel_history(request.auth_email)
    for entry in history:
        for k in ("requested_at", "approved_at", "start_date", "end_date"):
            v = entry.get(k)
            if hasattr(v, "strftime"):
                entry[k] = v.strftime("%Y-%m-%d %H:%M")
    return jsonify(history)


@app.route("/api/manager/travel-requests")
@login_required(roles=["manager"])
def manager_travel_requests():
    try:
        reqs = get_all_travel_requests()
        safe = []
        for r in reqs:
            if "_id" in r:
                r["_id"] = str(r["_id"])
            tm = r.get("travel_mode") or {}
            for k in ("requested_at", "approved_at", "start_date", "end_date"):
                v = tm.get(k)
                if hasattr(v, "strftime"):
                    tm[k] = v.strftime("%Y-%m-%d %H:%M")
            for k in ("created_at", "last_login", "last_active", "last_login_time"):
                v = r.get(k)
                if hasattr(v, "strftime"):
                    r[k] = v.strftime("%Y-%m-%d %H:%M")
            if tm and tm.get("status"):
                safe.append(r)
        return jsonify(safe)
    except Exception as e:
        import traceback
        print("[ERROR] manager_travel_requests: " + str(e))
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/api/manager/approve-travel", methods=["POST"])
@login_required(roles=["manager"])
def manager_approve_travel():
    data  = request.json or {}
    email = data.get("email", "").strip()
    if not email:
        return jsonify({"error": "Employee email required."}), 400
    approve_travel_mode(email, request.auth_email)
    log_activity(request.auth_email, "TRAVEL_MODE_APPROVED",
        "Approved travel mode for " + email, "Manager Dashboard", request.remote_addr, "LOW")
    emit_to_user(email, "travel_mode_resolved", {
        "status": "approved", "email": email,
        "message": "Your travel mode request has been approved by your manager."
    })
    return jsonify({"message": "Travel mode approved for " + email + "."})


@app.route("/api/manager/reject-travel", methods=["POST"])
@login_required(roles=["manager"])
def manager_reject_travel():
    data   = request.json or {}
    email  = data.get("email", "").strip()
    reason = data.get("reason", "Rejected by manager.")
    if not email:
        return jsonify({"error": "Employee email required."}), 400
    reject_travel_mode(email, request.auth_email)
    log_activity(request.auth_email, "TRAVEL_MODE_REJECTED",
        "Rejected travel mode for " + email + ". Reason: " + reason, "Manager Dashboard", request.remote_addr, "LOW")
    emit_to_user(email, "travel_mode_resolved", {
        "status": "rejected", "email": email,
        "message": "Your travel mode request was rejected. Reason: " + reason
    })
    return jsonify({"message": "Travel mode rejected for " + email + "."})


# ===============================================================
# MANAGER OWN TRAVEL ROUTES
# ===============================================================

@app.route("/api/manager/request-travel", methods=["POST"])
@login_required(roles=["manager"])
def manager_request_travel():
    data        = request.json or {}
    source      = data.get("source", "").strip()
    destination = data.get("destination", "").strip()
    start_date  = data.get("start_date", "").strip()
    end_date    = data.get("end_date", "").strip()
    reason      = data.get("reason", "").strip()

    if not source:
        return jsonify({"error": "Source location required."}), 400
    if not destination:
        return jsonify({"error": "Destination location required."}), 400
    if not start_date or not end_date:
        return jsonify({"error": "Start and end date/time required."}), 400
    if not reason:
        return jsonify({"error": "Reason for travel is required."}), 400

    try:
        start_dt = datetime.strptime(start_date, "%Y-%m-%dT%H:%M")
        end_dt   = datetime.strptime(end_date,   "%Y-%m-%dT%H:%M")
    except Exception:
        return jsonify({"error": "Invalid date format. Expected YYYY-MM-DDTHH:MM."}), 400
    if end_dt <= start_dt:
        return jsonify({"error": "End date/time must be after start."}), 400
    if (end_dt - start_dt).days > 30:
        return jsonify({"error": "Travel period cannot exceed 30 days."}), 400

    src_coords = data.get("src_coords")
    dst_coords = data.get("dst_coords")
    if src_coords:
        try:
            src_coords = {"lat": float(src_coords["lat"]), "lon": float(src_coords["lon"])}
        except Exception:
            src_coords = None
    if dst_coords:
        try:
            dst_coords = {"lat": float(dst_coords["lat"]), "lon": float(dst_coords["lon"])}
        except Exception:
            dst_coords = None

    request_travel_mode(
        request.auth_email, destination, start_dt, end_dt, reason,
        source=source, src_coords=src_coords, dst_coords=dst_coords
    )
    log_activity(
        request.auth_email, "MANAGER_TRAVEL_REQUESTED",
        "Manager travel requested: " + source + " -> " + destination + " (" + start_date + " to " + end_date + "). Reason: " + reason,
        "Manager Dashboard", request.remote_addr, "MEDIUM"
    )
    emit_to_admins("manager_travel_request", {
        "type":        "manager_travel_mode",
        "email":       request.auth_email,
        "source":      source,
        "destination": destination,
        "start_date":  start_date,
        "end_date":    end_date,
        "reason":      reason,
        "time":        datetime.utcnow().strftime("%H:%M:%S")
    })
    return jsonify({"message": "Travel request submitted to admin for approval."})


@app.route("/api/manager/travel-status")
@login_required(roles=["manager"])
def manager_travel_status():
    tm = get_travel_mode(request.auth_email)
    if not tm:
        return jsonify({})
    for k in ("requested_at", "approved_at", "start_date", "end_date"):
        v = tm.get(k)
        if hasattr(v, "strftime"):
            tm[k] = v.strftime("%Y-%m-%d %H:%M")
    return jsonify(tm)


@app.route("/api/manager/travel-history")
@login_required(roles=["manager"])
def manager_travel_history():
    history = get_travel_history(request.auth_email)
    for entry in history:
        for k in ("requested_at", "approved_at", "start_date", "end_date"):
            v = entry.get(k)
            if hasattr(v, "strftime"):
                entry[k] = v.strftime("%Y-%m-%d %H:%M")
    return jsonify(history)


# ===============================================================
# ADMIN — MANAGER TRAVEL APPROVAL ROUTES
# ===============================================================

@app.route("/api/admin/manager-travel-requests")
@login_required(roles=["admin"])
def admin_manager_travel_requests():
    try:
        docs = list(users_col.find(
            {"role": "manager", "travel_mode.status": {"$exists": True, "$ne": ""}},
            {"password": 0}
        ))
        safe = []
        for r in docs:
            if "_id" in r:
                r["_id"] = str(r["_id"])
            tm = r.get("travel_mode") or {}
            if not tm or not tm.get("status"):
                continue
            for k in ("requested_at", "approved_at", "start_date", "end_date"):
                v = tm.get(k)
                if hasattr(v, "strftime"):
                    tm[k] = v.strftime("%Y-%m-%d %H:%M")
            for k in ("created_at", "last_login", "last_active", "last_login_time", "device_enrolled_at"):
                v = r.get(k)
                if hasattr(v, "strftime"):
                    r[k] = v.strftime("%Y-%m-%d %H:%M")
            safe.append(r)
        print("[DEBUG] admin_manager_travel_requests returning " + str(len(safe)) + " records")
        return jsonify(safe)
    except Exception as e:
        import traceback
        print("[ERROR] admin_manager_travel_requests: " + str(e))
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/api/admin/approve-manager-travel", methods=["POST"])
@login_required(roles=["admin"])
def admin_approve_manager_travel():
    data  = request.json or {}
    email = data.get("email", "").strip()
    if not email:
        return jsonify({"error": "Manager email required."}), 400

    user = get_user_by_email(email)
    if not user or user.get("role") != "manager":
        return jsonify({"error": "User not found or not a manager."}), 404

    approve_travel_mode(email, request.auth_email)
    log_activity(
        request.auth_email, "ADMIN_TRAVEL_APPROVED",
        "Approved travel mode for manager " + email,
        "Admin Dashboard", request.remote_addr, "LOW"
    )
    emit_to_user(email, "travel_mode_resolved", {
        "status":  "approved",
        "email":   email,
        "message": "Your travel request has been approved by admin."
    })
    return jsonify({"message": "Travel mode approved for manager " + email + "."})


@app.route("/api/admin/reject-manager-travel", methods=["POST"])
@login_required(roles=["admin"])
def admin_reject_manager_travel():
    data   = request.json or {}
    email  = data.get("email", "").strip()
    reason = data.get("reason", "Rejected by admin.")
    if not email:
        return jsonify({"error": "Manager email required."}), 400

    user = get_user_by_email(email)
    if not user or user.get("role") != "manager":
        return jsonify({"error": "User not found or not a manager."}), 404

    reject_travel_mode(email, request.auth_email)
    log_activity(
        request.auth_email, "ADMIN_TRAVEL_REJECTED",
        "Rejected travel mode for manager " + email + ". Reason: " + reason,
        "Admin Dashboard", request.remote_addr, "LOW"
    )
    emit_to_user(email, "travel_mode_resolved", {
        "status":  "rejected",
        "email":   email,
        "message": "Your travel request was rejected by admin. Reason: " + reason
    })
    return jsonify({"message": "Travel mode rejected for manager " + email + "."})


# ===============================================================
# SESSION TIMEOUT
# ===============================================================

SESSION_TIMEOUT_SECONDS = 1800

@app.route("/api/auth/heartbeat", methods=["POST"])
@login_required(roles=["employee", "manager", "admin"])
def session_heartbeat():
    session["last_active"] = datetime.utcnow().isoformat()
    return jsonify({"status": "alive", "timeout": SESSION_TIMEOUT_SECONDS})


@app.route("/api/employee/request-extend", methods=["POST"])
@login_required(roles=["employee"])
def employee_request_extend():
    user = get_user_by_email(request.auth_email)
    _name = user["name"] if user else request.auth_email
    emit_to_managers("session_extend_request", {
        "email":   request.auth_email,
        "name":    _name,
        "time":    datetime.utcnow().strftime("%H:%M:%S"),
        "message": _name + " is requesting a session extension."
    })
    log_activity(request.auth_email, "SESSION_EXTEND_REQUESTED",
        "Employee requested session extension from admin.",
        "Employee Dashboard", request.remote_addr, "LOW")
    return jsonify({"message": "Extension request sent to admin."})


@app.route("/api/manager/approve-extend", methods=["POST"])
@app.route("/api/admin/approve-extend", methods=["POST"])
@login_required(roles=["manager", "admin"])
def manager_approve_extend():
    data  = request.json or {}
    email = data.get("email", "")
    if not email:
        return jsonify({"error": "Employee email required."}), 400
    emit_to_user(email, "session_extended", {
        "status":  "approved",
        "email":   email,
        "message": "Your session has been extended by your manager.",
        "extend_seconds": 1800
    })
    log_activity(request.auth_email, "SESSION_EXTEND_APPROVED",
        "Approved session extension for " + email,
        "Admin/Manager Dashboard", request.remote_addr, "LOW")
    return jsonify({"message": "Session extended for " + email + "."})


@app.route("/api/manager/reject-extend", methods=["POST"])
@app.route("/api/admin/reject-extend", methods=["POST"])
@login_required(roles=["manager", "admin"])
def manager_reject_extend():
    data  = request.json or {}
    email = data.get("email", "")
    if not email:
        return jsonify({"error": "Employee email required."}), 400
    emit_to_user(email, "session_extended", {
        "status":  "rejected",
        "email":   email,
        "message": "Your session extension was declined by your manager."
    })
    log_activity(request.auth_email, "SESSION_EXTEND_REJECTED",
        "Rejected session extension for " + email,
        "Admin/Manager Dashboard", request.remote_addr, "LOW")
    return jsonify({"message": "Session extension rejected for " + email + "."})


@app.route("/api/auth/check-session")
@login_required(roles=["employee", "manager", "admin"])
def check_session():
    last = session.get("last_active")
    if last:
        elapsed = (datetime.utcnow() - datetime.fromisoformat(last)).total_seconds()
        remaining = max(0, SESSION_TIMEOUT_SECONDS - elapsed)
        if elapsed > SESSION_TIMEOUT_SECONDS:
            return jsonify({"expired": True, "remaining": 0})
        return jsonify({"expired": False, "remaining": int(remaining)})
    session["last_active"] = datetime.utcnow().isoformat()
    return jsonify({"expired": False, "remaining": SESSION_TIMEOUT_SECONDS})


# ===============================================================
# ADMIN: SEND FILE TO USER(S)
# ===============================================================

@app.route("/api/admin/users-for-send", methods=["GET"])
@login_required(roles=["admin"])
def admin_users_for_send():
    role   = request.args.get("role", "all")
    search = request.args.get("q", "").strip().lower()
    users  = []
    if role in ("all", "manager"):
        users += get_all_users_by_role("manager")
    if role in ("all", "employee"):
        users += get_all_users_by_role("employee")
    if search:
        users = [u for u in users
                 if search in u.get("name","").lower()
                 or search in u.get("email","").lower()]
    return jsonify([
        {"name": u.get("name",""), "email": u.get("email",""), "role": u.get("role","")}
        for u in users if u.get("is_active", True)
    ])


@app.route("/api/admin/send-file", methods=["POST"])
@login_required(roles=["admin"])
def admin_send_file():
    if "file" not in request.files:
        return jsonify({"error": "No file provided."}), 400
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected."}), 400
    if not allowed_file(file.filename):
        return jsonify({"error": "File type not allowed."}), 400
    audience = request.form.get("audience", "specific")  # public | all_employees | all_managers | specific

    # Resolve recipients and visibility based on audience
    if audience == "public":
        recipients = []
        visibility  = "public"
    elif audience == "all_employees":
        recipients = [u["email"] for u in get_all_users_by_role("employee") if u.get("is_active", True)]
        visibility  = "private"
        if not recipients:
            return jsonify({"error": "No active employees found."}), 400
    elif audience == "all_managers":
        recipients = [u["email"] for u in get_all_users_by_role("manager") if u.get("is_active", True)]
        visibility  = "private"
        if not recipients:
            return jsonify({"error": "No active managers found."}), 400
    else:
        recipients_raw = request.form.get("recipients", "")
        recipients = [e.strip().lower() for e in recipients_raw.split(",") if e.strip()]
        visibility  = "private"
        if not recipients:
            return jsonify({"error": "Please select at least one recipient."}), 400

    ext         = file.filename.rsplit(".", 1)[1].lower()
    unique_name = uuid.uuid4().hex + "." + ext
    file_bytes  = file.read()
    file_size   = len(file_bytes)

    # Step 1: Save to GridFS FIRST — bytes must exist before metadata record
    gridfs_ok = save_file_to_gridfs(unique_name, file_bytes, file.content_type)
    if not gridfs_ok:
        return jsonify({"error": "Failed to store file in GridFS. Try again."}), 500

    # Step 2: Save metadata record — only after bytes are safely stored
    audience_label = {"public": "All (Public)", "all_employees": "All Employees",
                      "all_managers": "All Managers"}.get(audience, ", ".join(recipients))
    record = save_file_record(
        filename       = unique_name,
        original_name  = secure_filename(file.filename),
        file_size      = file_size,
        file_type      = file.content_type,
        visibility     = visibility,
        allowed_emails = recipients,
        uploaded_by    = request.auth_email,
        scan_clean     = None,
        scan_engine    = "admin-send",
        scan_detail    = "Sent by admin to: " + audience_label
    )

    # Step 3: DLP scan on temp disk copy (separate from GridFS bytes)
    save_path = os.path.join(UPLOAD_FOLDER, unique_name)
    try:
        with open(save_path, "wb") as _f:
            _f.write(file_bytes)
        from dlp.policy_engine import scan_and_enforce
        for recipient_email in recipients:
            scan_and_enforce(
                file_path         = save_path,
                user_email        = recipient_email,
                action_type       = "ADMIN_FILE_SEND",
                destination       = "admin→" + recipient_email,
                socketio_instance = socketio
            )
    except Exception as _dlp_err:
        print("[DLP] scan_and_enforce error (admin send): " + str(_dlp_err))
    finally:
        try:
            os.remove(save_path)
        except Exception:
            pass

    for email in recipients:
        emit_to_user(email, "file_uploaded", {
            "file_id":    record["_id"],
            "name":       record["original_name"],
            "visibility": visibility,
            "sent_by":    "Admin",
            "time":       datetime.utcnow().strftime("%H:%M:%S")
        })
    # For public files, broadcast to all connected users
    if audience == "public":
        socketio.emit("file_uploaded", {
            "file_id":    record["_id"],
            "name":       record["original_name"],
            "visibility": "public",
            "sent_by":    "Admin",
            "time":       datetime.utcnow().strftime("%H:%M:%S")
        })
    emit_to_managers("new_event", {
        "type":    "FILE_SENT_BY_ADMIN",
        "email":   request.auth_email,
        "user":    request.auth_email,
        "detail":  "Admin sent file: " + record["original_name"] + " to " + audience_label,
        "location":"Admin Dashboard",
        "risk":    "LOW",
        "blocked": False,
        "time":    datetime.utcnow().strftime("%H:%M:%S")
    })
    log_activity(request.auth_email, "ADMIN_FILE_SENT",
                 "Sent '" + record["original_name"] + "' to " + audience_label,
                 "Admin Dashboard", "internal", "LOW")
    msg = "File sent publicly to all users." if audience == "public" else \
          "File sent to " + str(len(recipients)) + " recipient(s)."
    return jsonify({
        "message":    msg,
        "file":       record,
        "recipients": recipients
    })


# ===============================================================
# VIRUS SCAN HELPER
# ===============================================================

def scan_file_for_virus(filepath):
    defender_path = r"C:\Program Files\Windows Defender\MpCmdRun.exe"
    if os.path.exists(defender_path):
        try:
            result = subprocess.run(
                [defender_path, "-Scan", "-ScanType", "3", "-File", filepath, "-DisableRemediation"],
                capture_output=True, text=True, timeout=30
            )
            output = result.stdout + result.stderr
            if "found no threats" in output.lower() or result.returncode == 0:
                return True, "Windows Defender", "No threats found"
            elif "threat" in output.lower() or result.returncode == 2:
                threat = "Unknown threat"
                for line in output.splitlines():
                    if "threat" in line.lower() and ":" in line:
                        threat = line.strip()
                        break
                return False, "Windows Defender", threat
        except Exception:
            pass

    try:
        result = subprocess.run(
            ["clamscan", "--no-summary", filepath],
            capture_output=True, text=True, timeout=30
        )
        output = result.stdout + result.stderr
        if result.returncode == 0:
            return True, "ClamAV", "No threats found"
        elif result.returncode == 1:
            for line in output.splitlines():
                if "FOUND" in line:
                    return False, "ClamAV", line.strip()
            return False, "ClamAV", "Malware detected"
    except FileNotFoundError:
        pass
    except Exception:
        pass

    return _heuristic_scan(filepath)


def _heuristic_scan(filepath):
    filename = os.path.basename(filepath).lower()

    parts = filename.split('.')
    if len(parts) > 2:
        last_ext = parts[-1]
        dangerous_exts = {'exe','bat','cmd','ps1','vbs','js','scr','com','pif','hta','jar','msi','reg'}
        if last_ext in dangerous_exts:
            return False, "Heuristic", "Dangerous double extension detected: ." + parts[-2] + "." + last_ext

    try:
        with open(filepath, 'rb') as f:
            header = f.read(512)

        exe_magic = b'\x4d\x5a'
        elf_magic = b'\x7fELF'

        if header.startswith(exe_magic):
            return False, "Heuristic", "File contains Windows executable (MZ) header - likely disguised EXE"
        if header.startswith(elf_magic):
            return False, "Heuristic", "File contains Linux ELF binary header"

        ext = os.path.splitext(filename)[1].lower()
        if ext in {'.docx', '.xlsx', '.pptx', '.doc', '.xls', '.ppt'}:
            with open(filepath, 'rb') as f:
                content = f.read(8192)
            macro_indicators = [b'VBA', b'AutoOpen', b'AutoExec', b'Shell(', b'CreateObject', b'WScript.Shell']
            for indicator in macro_indicators:
                if indicator in content:
                    return False, "Heuristic", "Macro/script indicator found in Office file: " + indicator.decode(errors='ignore')

        if ext in {'.txt', '.csv', '.json', '.xml', '.py'} and os.path.getsize(filepath) < 1_000_000:
            with open(filepath, 'r', errors='ignore') as f:
                text_content = f.read(10000).lower()
            danger_patterns = [
                'exec(base64', 'eval(base64', 'shellcode', 'payload',
                'reverse_shell', 'meterpreter', 'mimikatz',
                'invoke-expression', 'downloadstring', 'webclient',
            ]
            for pat in danger_patterns:
                if pat in text_content:
                    return False, "Heuristic", "Suspicious pattern found: " + pat

        return True, "Heuristic", "No threats detected (heuristic scan)"

    except Exception as e:
        return True, "Heuristic", "Scan skipped (error: " + str(e)[:50] + ")"


# ===============================================================
# EMPLOYEE FILE UPLOAD (with virus scan)
# ===============================================================

@app.route("/api/employee/upload-doc", methods=["POST"])
@login_required(roles=["employee"])
def employee_upload_doc():
    if "file" not in request.files:
        return jsonify({"error": "No file provided."}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected."}), 400
    if not allowed_file(file.filename):
        return jsonify({"error": "File type not allowed."}), 400

    ext         = file.filename.rsplit(".", 1)[1].lower()
    temp_name   = "scan_" + uuid.uuid4().hex + "." + ext
    temp_path   = os.path.join(UPLOAD_FOLDER, temp_name)
    file.save(temp_path)

    is_clean, engine, detail = scan_file_for_virus(temp_path)

    scan_result = {
        "clean":      is_clean,
        "engine":     engine,
        "detail":     detail,
        "file":       file.filename,
        "scanned_at": datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")
    }

    if not is_clean:
        try:
            os.remove(temp_path)
        except Exception:
            pass
        log_activity(
            request.auth_email, "VIRUS_DETECTED",
            "Virus detected in upload: " + file.filename + " - Engine: " + engine + " - " + detail,
            "Employee Upload", request.remote_addr, "HIGH"
        )
        log_security_event(
            request.auth_email, "VIRUS_UPLOAD_BLOCKED",
            "Blocked infected file: " + file.filename + " (" + engine + ": " + detail + ")",
            "Employee Upload", request.remote_addr, blocked=True
        )
        emit_to_managers("new_event", {
            "type":    "VIRUS_DETECTED",
            "email":   request.auth_email,
            "user":    request.auth_email,
            "detail":  "INFECTED FILE BLOCKED: " + file.filename + " - " + engine + ": " + detail,
            "location":"Employee Upload",
            "risk":    "HIGH",
            "blocked": True,
            "time":    datetime.utcnow().strftime("%H:%M:%S")
        })
        xai_virus_reason = (
            "The file \"" + file.filename + "\" was scanned by the " + engine + " engine "
            "and identified as malicious. Detail: " + detail + ". "
            "The file was deleted immediately and never stored on the server. "
            "This upload has been logged as a HIGH-risk security event."
        )
        return jsonify({
            "error":      "File rejected - virus/malware detected.",
            "scan":       scan_result,
            "xai_reason": xai_virus_reason
        }), 422

    final_name = uuid.uuid4().hex + "." + ext
    final_path = os.path.join(UPLOAD_FOLDER, final_name)
    os.rename(temp_path, final_path)
    file_size = os.path.getsize(final_path)

    # Save to GridFS then remove temp file
    with open(final_path, "rb") as _fh:
        final_bytes = _fh.read()
    save_file_to_gridfs(final_name, final_bytes, file.content_type)
    try:
        os.remove(final_path)
    except Exception:
        pass

    record = save_file_record(
        filename       = final_name,
        original_name  = secure_filename(file.filename),
        file_size      = file_size,
        file_type      = file.content_type,
        visibility     = "private",
        allowed_emails = [request.auth_email],
        uploaded_by    = request.auth_email,
        scan_clean     = True,
        scan_engine    = engine,
        scan_detail    = detail
    )

    log_activity(
        request.auth_email, "EMPLOYEE_FILE_UPLOADED",
        "Uploaded (clean): " + file.filename + " - Scan: " + engine,
        "Employee Upload", request.remote_addr, "LOW"
    )

    emit_to_managers("new_event", {
        "type":    "EMPLOYEE_FILE_UPLOADED",
        "email":   request.auth_email,
        "user":    request.auth_email,
        "detail":  "Employee uploaded file: " + file.filename + " - " + engine + ": clean",
        "location":"Employee Upload",
        "risk":    "LOW",
        "blocked": False,
        "time":    datetime.utcnow().strftime("%H:%M:%S")
    })

    return jsonify({
        "message": "File uploaded successfully - scan passed.",
        "file":    record,
        "scan":    scan_result
    })


# --- SocketIO -----------------------------------------------------------------

sid_to_email = {}
email_to_sids = {}

@socketio.on("connect")
def on_connect():
    print("[SOCKET] Client connected: " + request.sid)

@socketio.on("disconnect")
def on_disconnect():
    sid = request.sid
    email = sid_to_email.pop(sid, None)
    if email and email in email_to_sids:
        email_to_sids[email].discard(sid)
        if not email_to_sids[email]:
            del email_to_sids[email]
    print("[SOCKET] Client disconnected: " + sid)

@socketio.on("register")
def on_register(data):
    from flask_socketio import join_room
    email = data.get("email", "")
    role  = data.get("role", "")
    sid   = request.sid
    sid_to_email[sid] = email
    if email not in email_to_sids:
        email_to_sids[email] = set()
    email_to_sids[email].add(sid)
    if role == "manager" or role == "admin":
        join_room("managers")
    else:
        join_room("employees")
    join_room("user_" + email)
    print("[SOCKET] Registered " + email + " (" + role + ") sid=" + sid)

def emit_to_managers(event, data):
    socketio.emit(event, data)

def emit_to_employees(event, data):
    socketio.emit(event, data)

def emit_to_user(email, event, data):
    data["_target_email"] = email
    socketio.emit(event, data)

def emit_to_all(event, data):
    socketio.emit(event, data)

def emit_to_admins(event, data):
    socketio.emit(event, data)


# ===============================================================
# MODULE 2 - AI RISK SCORING ENGINE
# ===============================================================

import pickle
import numpy as np
from sklearn.preprocessing import StandardScaler

try:
    with open('model.pkl', 'rb') as f:
        model_data = pickle.load(f)
    ai_model        = model_data['model']
    ai_scaler       = model_data['scaler']
    ai_feature_cols = model_data['feature_cols']
    print("[AI] model.pkl loaded successfully.")
except Exception as e:
    ai_model = None
    print("[AI] Warning: Could not load model.pkl - " + str(e))


def build_employee_features(email):
    from models.user import logs_col, events_col

    logs   = list(logs_col.find({"email": email}))
    events = list(events_col.find({"email": email}))

    login_count       = sum(1 for l in logs if l.get('action') == 'LOGIN_SUCCESS')
    afterhours_logins = 0
    unique_pcs        = 1

    for l in logs:
        if l.get('action') == 'LOGIN_SUCCESS':
            ts = l.get('timestamp')
            if ts:
                hour = ts.hour if hasattr(ts, 'hour') else 0
                if hour < 8 or hour > 18:
                    afterhours_logins += 1

    file_access_count     = sum(1 for l in logs if l.get('action') in
                                ('FILE_VIEWED', 'FILE_DOWNLOADED', 'EMPLOYEE_FILE_UPLOADED'))
    removable_media_write = sum(1 for e in events if e.get('action') == 'USB_INSERTED')
    removable_media_read  = 0

    usb_attempts = sum(1 for e in events if e.get('action') == 'USB_INSERTED')

    email_count       = sum(1 for l in logs if 'EMAIL' in str(l.get('action', '')))
    total_attachments = sum(1 for l in logs if l.get('action') == 'ACCESS_REQUESTED_FORWARD')
    avg_email_size    = 0.0

    psych_risk = 50.0

    return [
        login_count, afterhours_logins, unique_pcs,
        file_access_count, removable_media_write, removable_media_read,
        usb_attempts, email_count, total_attachments,
        avg_email_size, psych_risk
    ]


@app.route("/api/manager/ai-risk-scores")
@login_required(roles=["manager", "admin"])
def ai_risk_scores():
    if not ai_model:
        return jsonify({"error": "AI model not loaded."}), 500

    employees = get_all_employees()
    results   = []

    for emp in employees:
        email = emp.get("email", "")
        name  = emp.get("name", "")
        prediction = 1  # default normal

        try:
            features = build_employee_features(email)
            X        = np.array([features])
            X_scaled = ai_scaler.transform(X)

            score      = ai_model.decision_function(X_scaled)[0]
            prediction = ai_model.predict(X_scaled)[0]

            risk_score = max(0, min(100, int((1 - score) * 50)))
            risk_level = "HIGH" if risk_score >= 70 else (
                         "MEDIUM" if risk_score >= 40 else "LOW")

            reasons = []
            f = features
            if f[1] > 0:  reasons.append(str(int(f[1])) + " after-hours login(s)")
            if f[6] > 0:  reasons.append(str(int(f[6])) + " USB attempt(s)")
            if f[3] > 10: reasons.append("High file access (" + str(int(f[3])) + " files)")
            if f[4] > 0:  reasons.append("Removable media write detected")
            if f[7] > 5:  reasons.append("High email activity (" + str(int(f[7])) + " emails)")
            if not reasons:
                reasons.append("Normal behavior detected")

            results.append({
                "email":      email,
                "name":       name,
                "risk_score": risk_score,
                "risk_level": risk_level,
                "reason":     ", ".join(reasons),
                "flagged":    bool(prediction == -1)
            })

        except Exception as e:
            results.append({
                "email":      email,
                "name":       name,
                "risk_score": 0,
                "risk_level": "LOW",
                "reason":     "Insufficient data",
                "flagged":    bool(prediction == -1)
            })

    results.sort(key=lambda x: x["risk_score"], reverse=True)
    return jsonify(results)


# ===============================================================
# PHONE DETECTION + FILE VIEWING STATE
# ===============================================================

# ── File viewing state — MongoDB-backed so agent polling survives restarts ──
_fv_col = _users_col_ref.database["file_viewing_state"]
try:
    _fv_col.create_index("email", unique=True)
except Exception:
    pass


class _FileViewingUsers:
    """dict-like {email: bool} backed by MongoDB."""

    def __init__(self):
        self._cache = {}
        for doc in _fv_col.find({}, {"_id": 0, "email": 1, "active": 1}):
            self._cache[doc["email"]] = doc.get("active", False)

    def __setitem__(self, email, active):
        self._cache[email] = bool(active)
        try:
            _fv_col.update_one(
                {"email": email},
                {"$set": {"email": email, "active": bool(active),
                           "updated_at": datetime.utcnow()}},
                upsert=True
            )
        except Exception:
            pass

    def get(self, email, default=False):
        return self._cache.get(email, default)

    def items(self):
        return self._cache.items()


file_viewing_users = _FileViewingUsers()

# Latest webcam JPEG per user — pushed by monitor.py, polled by browser
_cam_frames = {}

@app.route("/api/agent/push-frame", methods=["POST"])
def agent_push_frame():
    token = request.headers.get("X-Auth-Token") or request.args.get("token", "")
    if not token or token not in active_tokens:
        return jsonify({"error": "Unauthorized"}), 403
    email = active_tokens[token]
    if request.data:
        _cam_frames[email] = request.data
    return jsonify({"ok": True})


@app.route("/api/agent/camera-feed")
def agent_camera_feed():
    token = request.args.get("token", "")
    if not token or token not in active_tokens:
        return ("", 204)
    email = active_tokens[token]
    frame = _cam_frames.get(email)
    if not frame:
        # Also check if any user has a frame (manager viewing, employee agent running)
        if _cam_frames:
            frame = next(iter(_cam_frames.values()))
        else:
            return ("", 204)
    from flask import Response
    return Response(frame, mimetype="image/jpeg",
                    headers={"Cache-Control": "no-store, no-cache"})


@app.route("/api/agent/file-viewing", methods=["POST"])
@login_required(roles=["employee", "manager"])
def set_file_viewing():
    data   = request.json or {}
    active = data.get("active", False)
    token = request.headers.get("X-Auth-Token") or request.args.get("token", "")
    if token and token in active_tokens:
        email = active_tokens[token]
    elif "user_email" in session:
        email = session["user_email"]
    else:
        return jsonify({"error": "Unauthorized"}), 403
    file_viewing_users[email] = active
    print("[APP] File viewing set: " + email + " = " + str(active))
    return jsonify({"status": "ok", "email": email, "active": active})

@app.route("/api/agent/file-viewing-status")
def file_viewing_status():
    token = request.args.get("token", "")
    if not token or token not in active_tokens:
        return jsonify({"active": False})
    agent_email = active_tokens[token]
    # Check agent's own session first
    if file_viewing_users.get(agent_email, False):
        return jsonify({"active": True, "viewer": agent_email})
    # Check if ANY other logged-in user has a file open (e.g. manager)
    for em, active in file_viewing_users.items():
        if active:
            print("[APP] File status: agent=" + agent_email + " viewer=" + em)
            return jsonify({"active": True, "viewer": em})
    return jsonify({"active": False})


# In-memory phone detection flags {email: True} — short-lived, no need to persist
_phone_flags = {}


@app.route("/api/agent/phone-flag", methods=["POST"])
def agent_phone_flag():
    """Called by monitor.py when phone/recording is detected near screen."""
    token = request.headers.get("X-Auth-Token") or request.args.get("token", "")
    if not token or token not in active_tokens:
        return jsonify({"error": "Unauthorized"}), 403
    email = active_tokens[token]
    _phone_flags[email] = True
    # Immediately close file in browser via socket — no waiting for 2s poll
    emit_to_user(email, "new_event", {
        "type":   "CAMERA_BLOCKED",
        "detail": "Phone detected — file access suspended.",
        "time":   datetime.utcnow().strftime("%H:%M:%S")
    })
    log_security_event(email, "PHONE_DETECTED",
                       "Phone detected during file viewing — file access suspended.",
                       "Agent Monitor", "-", blocked=True)
    print("[APP] Phone flag set for: " + email)
    return jsonify({"ok": True})


@app.route("/api/agent/phone-status")
@login_required(roles=["employee", "manager"])
def phone_status():
    """Browser polls this every 2s; returns True once then clears the flag."""
    detected = _phone_flags.pop(request.auth_email, False)
    return jsonify({"phone_detected": bool(detected)})


# ===============================================================
# EMAIL + CLOUD UPLOAD ROUTES
# ===============================================================

import json
import hashlib

SENSITIVE_KEYWORDS = [
    'credit card', 'card number', 'cvv', 'ssn', 'social security',
    'password', 'passwd', 'api key', 'secret key', 'private key',
    'confidential', 'salary', 'payroll', 'bank account', 'routing number'
]

BLOCKED_FILE_TYPES = {
    'exe', 'bat', 'zip', 'rar', 'sql', 'db', 'key',
    'iso', 'dll', 'sh', 'cmd', 'vbs', 'msi'
}

ALLOWED_MIME_TYPES = {
    'application/pdf',
    'application/msword',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'application/vnd.ms-excel',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'text/plain', 'text/csv',
    'image/png', 'image/jpeg', 'image/gif',
    'application/vnd.ms-powerpoint',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
}

RESTRICTED_FOLDERS = {'Confidential', 'Private_Data', 'Sensitive_Projects'}


def get_role_limits():
    role = getattr(request, 'auth_role', 'employee')
    if role == 'manager':
        return {'external_email': True,  'max_attach': 5, 'max_upload_files': 5, 'max_upload_mb': 50,  'restricted_folders': True}
    elif role == 'intern':
        return {'external_email': False, 'max_attach': 1, 'max_upload_files': 1, 'max_upload_mb': 5,   'restricted_folders': False}
    else:
        return {'external_email': False, 'max_attach': 2, 'max_upload_files': 2, 'max_upload_mb': 10,  'restricted_folders': False}


def check_sensitive_keywords(text):
    t = text.lower()
    return [kw for kw in SENSITIVE_KEYWORDS if kw in t]


from pymongo import DESCENDING as _DESC
_mail_db = users_col.database
messages_col = _mail_db["messages"]
try:
    messages_col.create_index([("to", 1), ("sent_at", _DESC)])
    messages_col.create_index([("from_email", 1), ("sent_at", _DESC)])
except Exception:
    pass

def compute_file_hash(filepath):
    h = hashlib.sha256()
    try:
        with open(filepath, 'rb') as fh:
            for chunk in iter(lambda: fh.read(8192), b''):
                h.update(chunk)
    except Exception:
        pass
    return h.hexdigest()


def run_dlp_checks(filepath, filename):
    alerts, needs_ap = [], False
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

    if ext in BLOCKED_FILE_TYPES:
        alerts.append({'type': 'BLOCKED_EXTENSION', 'severity': 'CRITICAL',
                        'message': 'Blocked file type .' + ext + ' in "' + filename + '"'})
        needs_ap = True

    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as fh:
            content = fh.read(4096)
        found = check_sensitive_keywords(content)
        if found:
            alerts.append({'type': 'SENSITIVE_CONTENT', 'severity': 'HIGH',
                            'message': 'Sensitive keywords in "' + filename + '": ' + ", ".join(found)})
            needs_ap = True
    except Exception:
        pass

    return alerts, needs_ap


# ===============================================================
# INTERNAL MAIL ROUTES
# ===============================================================

@app.route('/api/mail/send', methods=['POST'])
@login_required(roles=['employee', 'manager'])
def mail_send():
    sender_email = request.auth_email
    sender_role  = request.auth_role

    # Parse: multipart/form-data (vfs_files attached) or plain JSON
    is_multipart = request.content_type and 'multipart/form-data' in request.content_type
    if is_multipart:
        import json as _json
        to_raw      = _json.loads(request.form.get('to', '[]'))
        subject     = request.form.get('subject', '').strip()
        body        = request.form.get('body', '').strip()
        thread_id   = request.form.get('thread_id', None)
        reply_to_id = request.form.get('reply_to_id', None)
        is_forward  = request.form.get('is_forward', 'false').lower() == 'true'
        vfs_files   = _json.loads(request.form.get('vfs_files', '[]'))
    else:
        data        = request.json or {}
        to_raw      = data.get('to', [])
        subject     = data.get('subject', '').strip()
        body        = data.get('body', '').strip()
        thread_id   = data.get('thread_id', None)
        reply_to_id = data.get('reply_to_id', None)
        is_forward  = data.get('is_forward', False)
        vfs_files   = []

    if not to_raw or not subject or not body:
        return jsonify({'success': False, 'error': 'Missing required fields: to, subject, body'}), 400

    to_list = [e.strip().lower() for e in (
        to_raw if isinstance(to_raw, list) else to_raw.split(',')
    ) if e.strip()]

    if not to_list:
        return jsonify({'success': False, 'error': 'No valid recipients'}), 400

    company_emails = {
        u['email'].strip().lower()
        for u in users_col.find({}, {'email': 1})
        if u.get('email')
    }
    not_found = [e for e in to_list if e not in company_emails]
    if not_found:
        return jsonify({'success': False,
                        'error': 'Not registered in system: ' + ", ".join(not_found)}), 400

    found_kw = check_sensitive_keywords(body)
    if found_kw and sender_role in ('employee', 'intern'):
        log_security_event(sender_email, 'EMAIL_SENSITIVE_BLOCKED',
                           'Blocked internal mail - keywords: ' + ", ".join(found_kw),
                           '-', '-', blocked=True)
        log_activity(sender_email, 'EMAIL_BLOCKED',
                     'Sensitive keywords: ' + ", ".join(found_kw) + ' | To: ' + ", ".join(to_list),
                     '-', '-', 'HIGH')
        xai_email_reason = (
            "Your email was blocked by the DLP engine because the message body contained "
            "sensitive keywords: " + ", ".join(found_kw) + ". "
            "Employees are not permitted to send messages containing passwords, card numbers, "
            "salary data, or other confidential terms via internal email."
        )
        return jsonify({'success': False, 'blocked': True,
                        'error': 'Message blocked - sensitive content detected: ' + ", ".join(found_kw),
                        'xai_reason': xai_email_reason}), 403

    if found_kw and sender_role == 'manager':
        log_security_event(sender_email, 'MANAGER_SENSITIVE_EMAIL',
                           'Sensitive keywords in mail: ' + ", ".join(found_kw),
                           '-', '-', blocked=False)

    # Resolve VFS file attachments — files already on server, referenced by fid
    if len(vfs_files) > 3:
        return jsonify({'success': False, 'error': 'Max 3 attachments allowed'}), 400

    attachments_meta = []
    for vf in vfs_files:
        fid  = str(vf.get('fid', ''))
        name = vf.get('name', 'unknown')
        if not fid:
            continue
        rec = get_file_by_id_unrestricted(fid)
        if not rec:
            return jsonify({'success': False, 'error': 'File not found in Files: ' + name}), 404
        # Grant recipients read access to this file
        _mail_db['files'].update_one(
            {'_id': rec['_id']},
            {'$addToSet': {'allowed_emails': {'$each': to_list}}}
        )
        attachments_meta.append({
            'file_id':       fid,
            'original_name': rec.get('original_name', name),
            'file_size':     rec.get('file_size', 0),
            'file_type':     rec.get('file_type', 'application/octet-stream'),
            'source':        'files'
        })

    now    = datetime.utcnow()
    msg_id = secrets.token_hex(16)
    if not thread_id:
        thread_id = msg_id

    messages_col.insert_one({
        '_id':          msg_id,
        'thread_id':    thread_id,
        'reply_to_id':  reply_to_id,
        'is_forward':   bool(is_forward),
        'from_email':   sender_email.lower(),
        'to':           to_list,
        'subject':      subject,
        'body':         body,
        'attachments':  attachments_meta,
        'sent_at':      now,
        'read_by':      [],
        'starred_by':   [],
        'deleted_for':  [],
        'sensitive':    bool(found_kw),
    })

    sender_user  = get_user_by_email(sender_email)
    sender_name  = sender_user.get('name', sender_email) if sender_user else sender_email
    attach_count = len(attachments_meta)

    for recipient in to_list:
        if recipient == sender_email.lower():
            continue
        emit_to_user(recipient, 'new_mail', {
            'msg_id':         msg_id,
            'thread_id':      thread_id,
            'from_email':     sender_email,
            'from_name':      sender_name,
            'subject':        subject,
            'preview':        body[:120] + ('...' if len(body) > 120 else ''),
            'sent_at':        now.strftime('%Y-%m-%d %H:%M'),
            'has_attachment': attach_count > 0,
        })

    attach_info = (' with ' + str(attach_count) + ' attachment(s)') if attach_count else ''
    log_activity(sender_email, 'EMAIL_SENT',
                 'Internal mail' + attach_info + ' -> ' + ", ".join(to_list) + ' | Subject: ' + subject[:60],
                 '-', '-', 'LOW')
    log_security_event(sender_email, 'EMAIL_SENT',
                       'Sent to ' + str(len(to_list)) + ' recipient(s)' + attach_info, '-', '-', blocked=False)

    return jsonify({
        'success':     True,
        'message':     'Message sent to ' + str(len(to_list)) + ' recipient(s)' + attach_info,
        'msg_id':      msg_id,
        'thread_id':   thread_id,
        'attachments': attachments_meta
    }), 200


@app.route('/api/mail/attachment/<file_id>', methods=['GET'])
@login_required(roles=['employee', 'manager'])
def mail_get_attachment(file_id):
    """Download a mail attachment — only accessible to sender / recipients."""
    from flask import send_file as _send_file
    email = request.auth_email.lower()
    rec   = get_file_by_id_unrestricted(file_id)
    if not rec:
        return jsonify({'error': 'Attachment not found'}), 404
    allowed = (
        rec.get('uploaded_by', '').lower() == email or
        email in [e.lower() for e in rec.get('allowed_emails', [])] or
        rec.get('visibility') == 'public'
    )
    if not allowed:
        return jsonify({'error': 'Access denied'}), 403
    file_bytes = get_file_from_gridfs(rec['filename'])
    if file_bytes is None:
        return jsonify({'error': 'File not found in storage'}), 404
    from flask import make_response, Response
    resp = make_response(Response(file_bytes, mimetype=rec.get('file_type', 'application/octet-stream')))
    resp.headers['Content-Disposition'] = 'attachment; filename="' + rec.get('original_name', rec['filename']) + '"'
    return resp


def _fmt_msg(d, viewer_email):
    return {
        'msg_id':       d['_id'],
        'thread_id':    d.get('thread_id', d['_id']),
        'reply_to_id':  d.get('reply_to_id'),
        'is_forward':   d.get('is_forward', False),
        'from_email':   d.get('from_email', ''),
        'to':           d.get('to', []),
        'subject':      d.get('subject', ''),
        'body':         d.get('body', ''),
        'sent_at':      d['sent_at'].strftime('%Y-%m-%d %H:%M') if hasattr(d.get('sent_at'), 'strftime') else '',
        'read':         viewer_email in d.get('read_by', []),
        'starred':      viewer_email in d.get('starred_by', []),
        'sensitive':    d.get('sensitive', False),
        'attachments':  d.get('attachments', []),
    }


@app.route('/api/mail/inbox', methods=['GET'])
@login_required(roles=['employee', 'manager'])
def mail_inbox():
    email  = request.auth_email.lower()
    q      = request.args.get('q', '').strip()
    query  = {'to': email, 'deleted_for': {'$nin': [email]}}
    if q:
        query['$or'] = [
            {'subject':    {'$regex': q, '$options': 'i'}},
            {'body':       {'$regex': q, '$options': 'i'}},
            {'from_email': {'$regex': q, '$options': 'i'}},
        ]
    docs = list(messages_col.find(query).sort('sent_at', _DESC).limit(200))
    return jsonify([_fmt_msg(d, email) for d in docs]), 200


@app.route('/api/mail/sent', methods=['GET'])
@login_required(roles=['employee', 'manager'])
def mail_sent():
    email  = request.auth_email.lower()
    q      = request.args.get('q', '').strip()
    query  = {'from_email': email, 'deleted_for': {'$nin': [email]}}
    if q:
        query['$or'] = [
            {'subject': {'$regex': q, '$options': 'i'}},
            {'body':    {'$regex': q, '$options': 'i'}},
        ]
    docs = list(messages_col.find(query).sort('sent_at', _DESC).limit(200))
    return jsonify([_fmt_msg(d, email) for d in docs]), 200


@app.route('/api/mail/starred', methods=['GET'])
@login_required(roles=['employee', 'manager'])
def mail_starred():
    email = request.auth_email.lower()
    docs  = list(messages_col.find(
        {'starred_by': email, 'deleted_for': {'$nin': [email]}}
    ).sort('sent_at', _DESC).limit(200))
    return jsonify([_fmt_msg(d, email) for d in docs]), 200


@app.route('/api/mail/thread/<thread_id>', methods=['GET'])
@login_required(roles=['employee', 'manager'])
def mail_thread(thread_id):
    email = request.auth_email.lower()
    docs  = list(messages_col.find(
        {'thread_id': thread_id,
         '$or': [{'to': email}, {'from_email': email}],
         'deleted_for': {'$nin': [email]}}
    ).sort('sent_at', 1))
    return jsonify([_fmt_msg(d, email) for d in docs]), 200


@app.route('/api/mail/read/<msg_id>', methods=['POST'])
@login_required(roles=['employee', 'manager'])
def mail_mark_read(msg_id):
    email = request.auth_email.lower()
    messages_col.update_one(
        {'_id': msg_id, 'to': email},
        {'$addToSet': {'read_by': email}}
    )
    return jsonify({'success': True}), 200


@app.route('/api/mail/star/<msg_id>', methods=['POST'])
@login_required(roles=['employee', 'manager'])
def mail_star(msg_id):
    email = request.auth_email.lower()
    msg   = messages_col.find_one({'_id': msg_id})
    if not msg:
        return jsonify({'success': False, 'error': 'Not found'}), 404
    starred = email in msg.get('starred_by', [])
    if starred:
        messages_col.update_one({'_id': msg_id}, {'$pull':     {'starred_by': email}})
    else:
        messages_col.update_one({'_id': msg_id}, {'$addToSet': {'starred_by': email}})
    return jsonify({'success': True, 'starred': not starred}), 200


@app.route('/api/mail/delete/<msg_id>', methods=['POST'])
@login_required(roles=['employee', 'manager'])
def mail_delete(msg_id):
    email = request.auth_email.lower()
    messages_col.update_one({'_id': msg_id}, {'$addToSet': {'deleted_for': email}})
    return jsonify({'success': True}), 200


@app.route('/api/mail/unread-count', methods=['GET'])
@login_required(roles=['employee', 'manager'])
def mail_unread_count():
    email = request.auth_email.lower()
    count = messages_col.count_documents({
        'to':          email,
        'read_by':     {'$nin': [email]},
        'deleted_for': {'$nin': [email]}
    })
    return jsonify({'unread': count}), 200


@app.route('/api/mail/users', methods=['GET'])
@login_required(roles=['employee', 'manager'])
def mail_users():
    docs   = list(users_col.find({}, {'_id': 0, 'name': 1, 'email': 1, 'role': 1}))
    result = []
    for d in docs:
        raw_email = d.get('email', '')
        if raw_email:
            result.append({
                'name':  d.get('name', ''),
                'email': raw_email.strip().lower(),
                'role':  d.get('role', '')
            })
    return jsonify(result), 200


@app.route('/api/mail/debug-users', methods=['GET'])
@login_required(roles=['employee', 'manager'])
def mail_debug_users():
    docs = list(users_col.find({}, {'_id': 0, 'name': 1, 'email': 1, 'role': 1}))
    return jsonify({'count': len(docs), 'users': docs}), 200


# ===============================================================
# CLOUD UPLOAD ROUTE
# ===============================================================

@app.route('/upload_to_drive', methods=['POST'])
@login_required(roles=['employee', 'manager'])
def upload_to_drive_route():
    email  = request.auth_email
    role   = request.auth_role
    limits = get_role_limits()
    folder = request.form.get('folder', 'Public_Docs')

    files     = request.files.getlist('files')
    vfs_files = json.loads(request.form.get('vfs_files', '[]'))
    total_files = len(files) + len(vfs_files)

    if total_files == 0:
        return jsonify({'success': False, 'error': 'No files provided'}), 400

    if total_files > limits['max_upload_files']:
        log_security_event(email, 'CLOUD_COUNT_VIOLATION',
                           str(total_files) + ' files attempted (limit ' + str(limits["max_upload_files"]) + ')',
                           '-', '-', blocked=True)
        return jsonify({
            'success': False, 'requires_approval': True,
            'alerts': [{'type': 'FILE_COUNT', 'severity': 'HIGH',
                        'message': 'Max ' + str(limits["max_upload_files"]) + ' files for role ' + role}]
        }), 403

    if folder in RESTRICTED_FOLDERS and not limits['restricted_folders']:
        log_security_event(email, 'CLOUD_FOLDER_VIOLATION',
                           'Role ' + role + ' tried restricted folder: ' + folder,
                           '-', '-', blocked=True)
        return jsonify({
            'success': False, 'blocked': True,
            'reason': 'Your role (' + role + ') cannot upload to restricted folder: ' + folder
        }), 403

    temp_paths, all_alerts = [], []
    requires_approval = False

    try:
        for f in files:
            fname = secure_filename(f.filename)
            temp  = os.path.join(UPLOAD_FOLDER, 'dlptmp_' + secrets.token_hex(8) + '_' + fname)
            f.save(temp)
            temp_paths.append((temp, fname))

        for temp, fname in temp_paths:
            size_mb = os.path.getsize(temp) / (1024 * 1024)
            if size_mb > limits['max_upload_mb']:
                all_alerts.append({'type': 'SIZE_EXCEEDED', 'severity': 'HIGH',
                                    'message': fname + ': ' + str(round(size_mb, 1)) + ' MB exceeds ' + str(limits["max_upload_mb"]) + ' MB'})
                requires_approval = True

            # ── Module 3 DLP scan ──────────────────────────────────────────────
            try:
                from dlp.policy_engine import scan_and_enforce
                dlp_result = scan_and_enforce(
                    file_path         = temp,
                    user_email        = email,
                    action_type       = "UPLOAD",
                    destination       = folder,
                    socketio_instance = socketio
                )
                if dlp_result and dlp_result.get("action") == "BLOCKED":
                    for temp2, _ in temp_paths:
                        try: os.remove(temp2)
                        except Exception: pass
                    return jsonify({
                        'success': False,
                        'blocked': True,
                        'reason':  dlp_result.get("reason", "DLP policy blocked this upload.")
                    }), 403
            except Exception as _dlp_err:
                print(f"[DLP] scan_and_enforce error (upload): {_dlp_err}")
            # ──────────────────────────────────────────────────────────────────

            alerts, needs_ap = run_dlp_checks(temp, fname)
            all_alerts.extend(alerts)
            if needs_ap:
                requires_approval = True

        for vf in vfs_files:
            vf_name = vf.get('name', 'unknown')
            ext = vf_name.rsplit('.', 1)[-1].lower() if '.' in vf_name else ''
            if ext in BLOCKED_FILE_TYPES:
                all_alerts.append({'type': 'BLOCKED_EXTENSION', 'severity': 'CRITICAL',
                                    'message': 'Blocked type .' + ext + ': "' + vf_name + '"'})
                requires_approval = True

        if requires_approval:
            upload_id = secrets.token_hex(16)
            log_security_event(email, 'CLOUD_UPLOAD_PENDING',
                               str(total_files) + ' file(s) pending review - folder: ' + folder,
                               '-', '-', blocked=False)
            emit_to_managers('upload_approval_request', {
                'upload_id':  upload_id,
                'user_email': email,
                'files':      [t[1] for t in temp_paths] + [v.get('name', '') for v in vfs_files],
                'folder':     folder,
                'alerts':     all_alerts,
                'timestamp':  datetime.now().isoformat(),
            })
            _xai_summary = '; '.join([a.get('message','') for a in all_alerts]) or 'DLP policy violation detected.'
            _xai_score   = 75 if any(a.get('severity') == 'CRITICAL' for a in all_alerts) else \
                           60 if any(a.get('severity') == 'HIGH'     for a in all_alerts) else 40
            _xai_level   = 'CRITICAL' if _xai_score >= 75 else 'HIGH' if _xai_score >= 60 else 'MEDIUM'
            _xai_explanation = [
                "This upload was flagged by the DLP policy engine and requires manager approval before it can proceed.",
                "Reason(s) detected: " + _xai_summary,
                "Risk level assigned: " + _xai_level + " (score " + str(_xai_score) + "/100). "
                "A manager must review and approve or deny this upload within 48 hours."
            ]
            return jsonify({
                'success':           False,
                'requires_approval': True,
                'upload_id':         upload_id,
                'alerts':            all_alerts,
                'message':           'Upload requires manager approval due to DLP violations',
                'xai': {
                    'risk_score':  _xai_score,
                    'risk_level':  _xai_level,
                    'summary':     _xai_summary,
                    'explanation': _xai_explanation
                }
            }), 200

        uploaded_names = [t[1] for t in temp_paths] + [v.get('name', '') for v in vfs_files]

        for fname in uploaded_names:
            log_activity(email, 'CLOUD_UPLOAD',
                         'Uploaded "' + fname + '" -> ' + folder, '-', '-', 'LOW')
            log_security_event(email, 'CLOUD_UPLOAD_SUCCESS',
                               '"' + fname + '" -> ' + folder, '-', '-', blocked=False)

        return jsonify({
            'success':        True,
            'uploaded_files': uploaded_names,
            'message':        str(len(uploaded_names)) + ' file(s) uploaded successfully to ' + folder
        }), 200

    except Exception as e:
        print('upload_to_drive error: ' + str(e))
        import traceback; traceback.print_exc()
        return jsonify({'success': False, 'error': str(e)}), 500

    finally:
        for temp, _ in temp_paths:
            try:
                os.remove(temp)
            except Exception:
                pass


# ===============================================================
# EMPLOYEE NOTIFICATIONS
# ===============================================================

@app.route("/api/employee/notifications")
@login_required(roles=["employee"])
def employee_notifications():
    return jsonify([])


# --- Run ----------------------------------------------------------------------
  if __name__ == "__main__":
    print("=" * 60)
    print("  XAI-ITD-DLP Framework - Module 1 Starting...")
    print("=" * 60)
    try:
        normalize_existing_emails()
    except Exception as e:
        print("[APP] normalize_existing_emails error: " + str(e))
    print("[APP] Visit: http://127.0.0.1:5000")
    from ml.scheduler import start_scheduler
    start_scheduler(socketio_instance=socketio)
    socketio.run(app, host="0.0.0.0", port=5000, debug=False, use_reloader=False)
